# tools/pdfto/pdf_to_ppt_pptx_converter.py
import os
# import pdf2pptx

# def convert_pdf_to_pptx(input_file, output_file):
#     pdf2pptx(input_file, output_file)

# def clean_temp_files(file_path):
#     if os.path.exists(file_path):
#         os.remove(file_path)

import os
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches
from pdf2image import convert_from_bytes
from wand.image import Image as WandImage
from PIL import Image
from django.http import FileResponse
from PyPDF2 import PdfReader

# from reportlab.lib.pagesizes import letter
# from reportlab.platypus import SimpleDocTemplate, Paragraph
# from reportlab.lib.styles import getSampleStyleSheet
# from reportlab.lib.units import inch as Inches

# import os
import zipfile
# from PyPDF2 import PdfReader
# from io import BytesIO
# from django.http import FileResponse
# from pptx import Presentation
# from PIL import Image

# def convert_pdf_to_ppt_pptx(pdf_content, output_folder):
#     ppt_paths = []

#     with BytesIO(pdf_content) as pdf_buffer:
#         pdf_reader = PdfReader(pdf_buffer)
#         total_pages = len(pdf_reader.pages)

#         for page_number in range(total_pages):
#             pdf_page = pdf_reader.pages[page_number]
#             pdf_bytes = BytesIO()
#             pdf_page.write_to_stream(pdf_bytes, encryption_key=None)
#             pdf_bytes.seek(0)

#             try:
#                 # Convert PDF page to PNG
#                 with WandImage(file=pdf_bytes) as img:
#                     img.format = 'png'  # Specify desired image format
#                     png_bytes = img.make_blob()

#                 # Save PNG file
#                 png_path = os.path.join(output_folder, f"page_{page_number + 1}.png")
#                 with open(png_path, 'wb') as f:
#                     f.write(png_bytes)

#                 # Create a PowerPoint slide with the image
#                 ppt_path = create_ppt_slide(png_path, output_folder, page_number + 1)
#                 ppt_paths.append(ppt_path)

#             except Exception as e:
#                 print(f"Error converting PDF page to image: {e}")

#     return ppt_paths


# def create_ppt_slide(image_path, output_folder, slide_number):
#     output_path = os.path.join(output_folder, f"slide_{slide_number}.pptx")

#     prs = Presentation()
#     blank_slide_layout = prs.slide_layouts[6]  # Use a blank layout for the slide
#     slide = prs.slides.add_slide(blank_slide_layout)
    
#     # Use Inches correctly
#     left = top = Inches(1)
    
#     # Add the image to the slide
#     slide.shapes.add_picture(image_path, left, top)

#     # Save the PowerPoint slide
#     prs.save(output_path)

#     return output_path


def create_zip_archives(file_paths, output_zip):
    with zipfile.ZipFile(output_zip, 'w') as zipf:
        for file_path in file_paths:
            zipf.write(file_path, os.path.basename(file_path))

    return output_zip


def convert_pdf_to_ppt_pptx(pdf_content, output_folder):
    ppt_paths = []

    with BytesIO(pdf_content) as pdf_buffer:
        pdf_reader = PdfReader(pdf_buffer)
        total_pages = len(pdf_reader.pages)

        for page_number in range(total_pages):
            pdf_page = pdf_reader.pages[page_number]
            pdf_bytes = BytesIO()
            pdf_page.write_to_stream(pdf_bytes, encryption_key=None)
            pdf_bytes.seek(0)

            try:
                # Convert PDF page to PNG
                with WandImage(file=pdf_bytes) as img:
                    img.format = 'png'  # Specify desired image format
                    png_bytes = img.make_blob()

                # Save PNG file
                png_path = os.path.join(output_folder, f"page_{page_number + 1}.png")
                with open(png_path, 'wb') as f:
                    f.write(png_bytes)

                # Create a PowerPoint slide with the image
                ppt_path = create_ppt_slide(png_path, output_folder, page_number + 1)
                ppt_paths.append(ppt_path)

            except Exception as e:
                print(f"Error converting PDF page to image: {e}")

    return ppt_paths


def create_ppt_slide(image_path, output_folder, slide_number):
    output_path = os.path.join(output_folder, f"slide_{slide_number}.pptx")

    prs = Presentation()
    blank_slide_layout = prs.slide_layouts[6]  # Use a blank layout for the slide
    slide = prs.slides.add_slide(blank_slide_layout)
    
    # Use Inches correctly
    left = top = Inches * 1
    
    # Add the image to the slide
    slide.shapes.add_picture(image_path, left, top)

    # Save the PowerPoint slide
    prs.save(output_path)

    return output_path



# def convert_pdf_to_ppt_pptx(pdf_content, output_folder):
#     ppt_paths = []

#     with BytesIO(pdf_content) as pdf_buffer:
#         pdf_reader = PdfReader(pdf_buffer)
#         total_pages = pdf_reader.pages
#         print(f'total pages : {total_pages}')

#         for page_number in range(total_pages):
#             # Convert PDF page to an image
#             pdf_page = pdf_reader.getPage(page_number)
#             pdf_bytes = BytesIO()
#             pdf_page.writeTo(pdf_bytes)
#             pdf_bytes.seek(0)
#             image = Image.open(pdf_bytes)
#             image_path = os.path.join(output_folder, f"page_{page_number + 1}.png")
#             image.save(image_path)

#             # Create a PowerPoint slide with the image
#             ppt_path = create_ppt_slide(image_path, output_folder, page_number + 1)
#             ppt_paths.append(ppt_path)

#     return ppt_paths

# def convert_pdf_to_ppt_pptx(pdf_content, output_folder):
#     ppt_paths = []

#     with BytesIO(pdf_content) as pdf_buffer:
#         pdf_reader = PdfReader(pdf_buffer)
#         total_pages = len(pdf_reader.pages)  # Corrected line

#         for page_number in range(total_pages):
#             # Convert PDF page to an image
#             pdf_page = pdf_reader.pages[page_number]
#             pdf_bytes = BytesIO()
#             pdf_page.write_to_stream(pdf_bytes, encryption_key=None)
#             pdf_bytes.seek(0)
            
#         #    # Create a new BytesIO object and write pdf_bytes into it
#         #     image_bytes = BytesIO()
#         #     image_bytes.write(pdf_bytes.getvalue())
#         #     image_bytes.seek(0)

#         #     # Open the image using PIL
#         #     image = Image.open(image_bytes)
#             image_path = os.path.join(output_folder, f"page_{page_number + 1}.png")
#         #     image.save(image_path)

#             try:
#                 with WandImage(file=pdf_bytes) as img:
#                     img.format = 'png'  # Specify desired image format
#                     img.save(filename=image_path)
#             except Exception as e:
#                 # Handle cases where image conversion might fail
#                 print(f"Error converting PDF page to image: {e}")
#                 # Handle the error appropriately, e.g., log it or raise a custom exception


#             # Create a PowerPoint slide with the image
#             ppt_path = create_ppt_slide(image_path, output_folder, page_number + 1)
#             ppt_paths.append(ppt_path)

#     return ppt_paths


# def create_ppt_slide(image_path, output_folder, slide_number):
#     output_path = os.path.join(output_folder, f"slide_{slide_number}.pptx")

#     prs = Presentation()
#     blank_slide_layout = prs.slide_layouts[6]  # Use a blank layout for the slide

#     slide = prs.slides.add_slide(blank_slide_layout)
#     # left = top = Inches(1)
#     left = top = Inches * 1
#     slide.shapes.add_picture(image_path, left, top)

#     prs.save(output_path)

#     return output_path


def clean_up_temp_files(file_paths):
    for file_path in file_paths:
        os.remove(file_path)

# def create_zip_archives(file_paths, zip_file_path):
#     """
#     Create a zip archive containing files specified by their paths.

#     :param file_paths: List of file paths to be included in the zip archive.
#     :param zip_file_path: Path to the output zip file.
#     """
#     with zipfile.ZipFile(zip_file_path, 'w') as zipf:
#         for file_path in file_paths:
#             zipf.write(file_path, os.path.basename(file_path))

# def convert_pdf_to_ppt_pptx(pdf_content, output_folder):
#     pdf_reader = PdfReader(pdf_content)

#     ppt_paths = []
#     page_size = (pdf_reader.pages[0].mediabox.width, pdf_reader.pages[0].mediabox.height)

#     for page_number, pdf_page in enumerate(pdf_reader.pages):
#         # Extract content from the PDF page
#         pdf_text = pdf_page.extract_text()
#         # Create a PowerPoint slide for each page
#         ppt_path = create_ppt_slide(pdf_text, output_folder, page_number + 1, page_size)
#         ppt_paths.append(ppt_path)

#     return ppt_paths

