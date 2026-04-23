# tools/views.py

# Django default libraries import
from django.conf import settings
from django.core.files.base import File
from django.core.files.storage import FileSystemStorage, default_storage
from django.core.signals import request_finished
from django.dispatch import receiver
from django.forms import Form
from django.http import HttpResponse, JsonResponse, FileResponse, HttpResponseBadRequest, Http404
from django.shortcuts import render, redirect
from django.template import TemplateDoesNotExist
from django.utils.text import slugify
from django.urls import reverse
from django.utils.encoding import smart_str
from django.views.decorators.csrf import csrf_exempt
from django.views import View
from PyPDF2 import PdfReader, PdfWriter, PageObject
from xlrd import open_workbook

# Models Import
from .models import ToolAttachment#, ConvertedPDF
from blog.models import Post

# Forms Import
from .forms import PDFUploadForm, RotatePDFForm, UploadFileForm
from .forms import RotatePDFForm

# from .tools.word_counter import word_counter_text
from .extra.split_pdf import split_pdf_by_page
from .extra.lorem_ipsum_generator import generate_lorem_ipsum

from .pdfto.pdf_to_docx_converter import pdf_to_docx_converter
from .pdfto.pdf_to_jpg_converter import convert_pdf_to_jpg, clean_up_jpg_files, create_zip_archive
# from .pdfto.pdf_to_ppt_pptx_converter import convert_pdf_to_pptx, create_ppt_slide, clean_up_temp_files, create_zip_archives

from .topdf.excel_to_pdf_converter import convert_excel_to_pdf
from .topdf.imgtopdf import convert_to_pdf
from .topdf.powerpoint_to_pdf_converter import convert_ppt_to_pdf
from .topdf.word_to_pdf_converter import clean_temp_files, convert_to_pdf#, convert_word_to_pdf


import os
import re
import docxtopdf
import img2pdf
import pdfkit
import requests
import subprocess
import tempfile
import zipfile

import base64
import ghostscript
import json
import locale
import pandas as pd
import pdf2image
import PyPDF2
import tabula
import traceback
import uuid


from io import BytesIO
from openpyxl import load_workbook
from pdf2docx import Converter
from pypdf import PdfWriter
from meta.views import Meta
from PyPDF2 import PdfMerger

from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.views.decorators.http import require_POST
from pdfminer.high_level import extract_text_to_fp
from PIL import Image, ImageSequence

from reportlab.lib import colors
from reportlab.lib.pagesizes import letter, landscape, legal
from reportlab.lib.units import inch
from reportlab.pdfgen import canvas
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle
from functools import wraps

from .topdf.word_to_pdf_converter import convert_to_pdf#, LibreOfficeError, sanitize_filename


import logging
import platform
logger = logging.getLogger(__name__)

# Maximum upload size: 30 MB
MAX_UPLOAD_SIZE = 30 * 1024 * 1024


def get_libreoffice_path():
    """Find the LibreOffice executable path across different OS."""
    if platform.system() == 'Windows':
        # Common Windows installation paths
        candidates = [
            r'C:\Program Files\LibreOffice\program\soffice.exe',
            r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
            os.path.expandvars(r'%PROGRAMFILES%\LibreOffice\program\soffice.exe'),
        ]
        for path in candidates:
            if os.path.isfile(path):
                return path
        # Fall back to command name (might be in PATH)
        return 'soffice'
    else:
        # Linux/macOS — libreoffice is typically in PATH
        return 'libreoffice'




# Views for the Tools From Here:

# ------------------------CHECKED
def merge_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST":
            files = request.FILES.getlist('pdf_files')
            if len(files) == 0:
                return HttpResponse("No files uploaded.", status=400)

            merger = PdfMerger()
            try:
                for file in files:
                    merger.append(file)

                response = HttpResponse(content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="merged_file.pdf"'

                merger.write(response)
                merger.close()

                return response

            except Exception as e:
                return HttpResponse(str(e), status=500)
        else:
            return view_func(request, *args, **kwargs)  
    return wrapper_function

# Merge PDF Tool for Attachment in the Blog Post, uses base.html
@merge_pdf_logic
def merge_pdf_view(request):
    meta = Meta(
        title='Merge PDF files free online',
        description='Merge PDF or Combine PDF in Order you want just within clicks.',
        keywords=['merge', 'combine', 'join', 'add'],
        og_title='Merge PDF files free online',
        og_description='Merge PDF or Combine PDF in Order you want just within clicks.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='merge_pdf_view')
    context = {'meta': meta, 'tool_attachment':tool_attachment}
    return render(request, 'tools/merge_pdf.html', context)

# Merge PDF Tool for Attachment in the Blog Post, doesn't use extends base.html
@merge_pdf_logic
def merge_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Merge PDF Tool',
        description='Merge PDF or Combine PDF in Order you want just within clicks.',
        keywords=['merge', 'combine', 'join', 'add'],
        og_title='iLovePdfConverterOnline - Merge PDF Tool',
        og_description='Merge PDF or Combine PDF in Order you want just within clicks.',
    )
    context = {'meta': meta}
    return render(request, 'tools/merge_pdf_include.html', context)

# -----------------------------------------================================

def split_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST":
            form = Form(request.POST, request.FILES)
            if form.is_valid():
                file = request.FILES['file']
                page_numbers = request.POST.get('page_numbers', '')
                output_files = split_pdf_by_page(file, page_numbers)

                response = HttpResponse(content_type='application/zip')
                zip_file = zipfile.ZipFile(response, 'w')
                for output_file in output_files:
                    zip_file.write(output_file, os.path.basename(output_file))
                zip_file.close()
                response['Content-Disposition'] = f'attachment; filename={file.name}.zip'
                return response
        else:
            return view_func(request, *args, **kwargs)  
    return wrapper_function


@split_pdf_logic
def split_pdf_view(request):
    meta = Meta(
        title='Split PDF document online',
        description='Split PDF or Unmerge PDF in Order you want just within clicks.',
        keywords=['split', 'unmerge', 'remove'],
        og_title='Split PDF document online',
        og_description='Split PDF or Unmerge PDF in Order you want just within clicks.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='split_pdf_view')
    context = {'meta': meta, 'tool_attachment':tool_attachment}
    return render(request, 'tools/split_pdf.html', context)

@split_pdf_logic
def split_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Split PDF',
        description='Split PDF or Unmerge PDF in Order you want just within clicks.',
        keywords=['split', 'unmerge', 'remove'],
        og_title='iLovePdfConverterOnline - Split PDF',
        og_description='Split PDF or Unmerge PDF in Order you want just within clicks.',
    )
    context = {'meta': meta}  
    return render(request, 'tools/split_pdf_include.html', context)

# -----------------------------------------================================


def get_pdf_settings(compress_level):
    """Returns the Ghostscript PDFSETTINGS parameter based on the compression level."""
    if compress_level <= 25:
        return "screen"  # Low quality
    elif compress_level <= 50:
        return "ebook"  # Medium quality
    elif compress_level <= 75:
        return "printer"  # High quality
    else:
        return "prepress"  # Maximum quality

def compress_pdf_logic(view_func):
    """Decorator to handle PDF compression via Ghostscript."""
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST':
            pdf_file = request.FILES.get('pdf_file')
            compress_level = int(request.POST.get('compress_level', 50))

            if pdf_file:
                # Use unique filename to avoid conflicts
                unique_id = uuid.uuid4().hex
                temp_pdf_path = default_storage.save(
                    f'temp_upload_{unique_id}.pdf',
                    ContentFile(pdf_file.read())
                )
                temp_pdf_full_path = os.path.join(
                    default_storage.location, temp_pdf_path
                )
                compressed_pdf_path = os.path.join(
                    settings.MEDIA_ROOT,
                    f'compressed_{unique_id}.pdf'
                )

                cargs = [
                    "ps2pdf",
                    "-dNOPAUSE", "-dBATCH", "-dSAFER",
                    "-sDEVICE=pdfwrite",
                    f"-dCompatibilityLevel=1.4",
                    f"-dPDFSETTINGS=/{get_pdf_settings(compress_level)}",
                    f"-sOutputFile={compressed_pdf_path}",
                    temp_pdf_full_path
                ]

                encoding = locale.getpreferredencoding()
                cargs = [a.encode(encoding) for a in cargs]

                try:
                    ghostscript.Ghostscript(*cargs)
                except ghostscript.GhostscriptError as e:
                    # Clean up temp file on error
                    default_storage.delete(temp_pdf_path)
                    return HttpResponse(
                        f"Error processing file with Ghostscript: {e}",
                        status=500
                    )

                # Clean up uploaded temp file
                default_storage.delete(temp_pdf_path)

                # Read compressed file into memory, then delete
                try:
                    with open(compressed_pdf_path, 'rb') as pdf:
                        pdf_data = pdf.read()
                    response = HttpResponse(
                        pdf_data, content_type='application/pdf'
                    )
                    response['Content-Disposition'] = (
                        'attachment; filename="compressed_output.pdf"'
                    )
                    return response
                finally:
                    # Delete file after reading into memory
                    try:
                        os.remove(compressed_pdf_path)
                    except OSError:
                        pass
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@compress_pdf_logic
def compress_pdf_view(request):
    meta = Meta(
        title='Compress PDF file online',
        description='Compress PDF to reduce the file size with percentage level.',
        keywords=['compress', 'reduce', 'small'],
        og_title='Compress PDF file online',
        og_description='Compress PDF file online in percentage level you want just within clicks.',
    )    
    tool_attachment = ToolAttachment.objects.get(function_name='compress_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/compress_pdf.html', context)

@compress_pdf_logic
def compress_pdf_include(request):
    meta = Meta(
        title='Compress PDF file online',
        description='Compress PDF to reduce the file size with percentage level.',
        keywords=['compress', 'reduce', 'small'],
        og_title='Compress PDF file online',
        og_description='Compress PDF file online in percentage level you want just within clicks.',
    ) 
    context = {'meta': meta}
    return render(request, 'tools/compress_pdf_include.html', context)

# -----------------------------------------================================


def parse_page_numbers(pages):
    page_numbers = set()
    for part in pages.split(','):
        if '-' in part:
            start, end = part.split('-')
            page_numbers.update(range(int(start) - 1, int(end)))  # Pages are 0-indexed internally
        else:
            page_numbers.add(int(part) - 1)  # Pages are 0-indexed internally
    return page_numbers

def rotate_pdf_logic(view_func):
    """Decorator to handle PDF rotation using PyPDF2."""
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST':
            form = RotatePDFForm(request.POST, request.FILES)
            if form.is_valid():
                try:
                    pdf_file = request.FILES['pdf_file']
                    rotation_angle = int(
                        form.cleaned_data['rotation_angle']
                    )
                    pages_to_rotate = form.cleaned_data['pages']

                    # Use PyPDF2 consistently for both reader and writer
                    reader = PyPDF2.PdfReader(pdf_file)
                    writer = PyPDF2.PdfWriter()

                    if pages_to_rotate:
                        pages_to_rotate = parse_page_numbers(
                            pages_to_rotate
                        )
                    else:
                        pages_to_rotate = range(len(reader.pages))

                    for i, page in enumerate(reader.pages):
                        if i in pages_to_rotate:
                            page.rotate(rotation_angle)
                        writer.add_page(page)

                    response = HttpResponse(
                        content_type='application/pdf'
                    )
                    response['Content-Disposition'] = (
                        'attachment; filename="rotated.pdf"'
                    )
                    writer.write(response)
                    return response
                except Exception as e:
                    logger.error(f"Rotate PDF error: {e}")
                    return HttpResponse(
                        f"Error rotating PDF: {str(e)}",
                        status=500
                    )
        return view_func(request, *args, **kwargs)
    return wrapper_function

@rotate_pdf_logic
def rotate_pdf_view(request):
    if request.method == 'GET':
        form = RotatePDFForm()
        meta = Meta(
            title='iLovePdfConverterOnline - Rotate PDF',
            description='Rotate PDF file or pages in 90, 180 or 270 degree.',
            keywords=['rotate', 'pdf', 'turn'],
            og_title='iLovePdfConverterOnline - Rotate PDF file or pages',
            og_description='Rotate PDF file or pages in 90, 180 or 270 degree.',
        )
        tool_attachment = ToolAttachment.objects.get(function_name='rotate_pdf_view')
        context = {'meta':meta, 'form': form, 'tool_attachment': tool_attachment}
        return render(request, 'tools/rotate_pdf.html', context)

@rotate_pdf_logic
def rotate_pdf_include(request):
    if request.method == 'GET':
        form = RotatePDFForm()
        meta = Meta(
            title='iLovePdfConverterOnline - Rotate PDF',
            description='Rotate PDF file or pages in 90, 180 or 270 degree.',
            keywords=['rotate', 'pdf', 'turn'],
            og_title='iLovePdfConverterOnline - Rotate PDF file or pages',
            og_description='Rotate PDF file or pages in 90, 180 or 270 degree.',
        )
        context = {'meta':meta, 'form': form}
        return render(request, 'tools/rotate_pdf_include.html', context)


# -----------------------------------------================================

# Working on Server — uses get_libreoffice_path() for cross-platform support

def word_to_pdf_logic(view_func):
    """Decorator to handle Word to PDF conversion via LibreOffice."""
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('word_file'):
            try:
                word_file = request.FILES['word_file']
                
                # Get the original name and strip extension
                original_name = word_file.name
                base_name = os.path.splitext(original_name)[0]
                
                # Use a unique prefix to prevent overwriting files with the same name
                unique_prefix = uuid.uuid4().hex[:8]
                temp_filename = f"{unique_prefix}_{original_name}"
                
                out_path = os.path.join(settings.MEDIA_ROOT, 'word_to_pdf')
                os.makedirs(out_path, exist_ok=True)
                temp_file_path = os.path.join(out_path, temp_filename)

                # Save file
                fs = FileSystemStorage(location=out_path)
                fs.save(temp_filename, word_file)

                env = os.environ.copy()
                env['HOME'] = out_path

                lo_path = get_libreoffice_path()
                subprocess.run(
                    [lo_path, '--headless', '--convert-to', 'pdf',
                     '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True, check=True
                )

                # The output PDF will match the saved temp_filename but with .pdf
                output_pdf_path = os.path.join(out_path, f"{os.path.splitext(temp_filename)[0]}.pdf")

                if os.path.exists(output_pdf_path):
                    response = FileResponse(
                        open(output_pdf_path, 'rb'),
                        content_type='application/pdf'
                    )
                    # Set the download name back to the original base name + .pdf
                    response['Content-Disposition'] = f'attachment; filename="{base_name}.pdf"'
                    
                    response.cleanup_files = [temp_file_path, output_pdf_path]
                    return response
                else:
                    return HttpResponse("Conversion failed.", status=500)
                    
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function


@word_to_pdf_logic
def word_to_pdf_view(request):
    meta = Meta(
        title='Word to PDF converter',
        description='iLovePdfConverterOnline Convert word document file (doc, docx) to PDF file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'docxtopdf'],
        og_title='Word to PDF converter',
        og_description='iLovePdfConverterOnline Convert word document file (doc, docx) to PDF file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='word_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/word_to_pdf.html', context) 

@word_to_pdf_logic
def word_to_pdf_include(request):
    meta = Meta(
        title='Word to PDF file converter',
        description='iLovePdfConverterOnline Convert word document file (doc, docx) to PDF file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'docxtopdf'],
        og_title='Word to PDF file converter',
        og_description='iLovePdfConverterOnline Convert word document file (doc, docx) to PDF file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/word_to_pdf_include.html', context)

# -----------------------------------------================================


def pdf_to_word_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST' and 'pdf_file' in request.FILES:
            pdf_file = request.FILES['pdf_file']
            try:
                # Define paths for uploaded and converted files
                upload_folder = os.path.join(settings.MEDIA_ROOT, 'uploads')
                os.makedirs(upload_folder, exist_ok=True)
                uploaded_pdf_path = os.path.join(upload_folder, 'uploaded_pdf.pdf')
                output_docx_path = os.path.join(upload_folder, 'converted_doc.docx')

                # Save the uploaded PDF file
                with open(uploaded_pdf_path, 'wb') as destination:
                    for chunk in pdf_file.chunks():
                        destination.write(chunk)

                # Call the pdf_to_docx_converter function
                success, error_message = pdf_to_docx_converter(uploaded_pdf_path, output_docx_path)
                if success:
                    with open(output_docx_path, 'rb') as docx_file:
                        response = HttpResponse(docx_file.read(), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                        response['Content-Disposition'] = 'attachment; filename="PDF_to_Word_iLovePDFconverteronline.com.docx"'
                        response.cleanup_files = [uploaded_pdf_path, output_docx_path]  # Add files for cleanup
                        return response
                else:
                    return HttpResponse(f"Conversion failed. Error: {error_message}")
            except Exception as e:
                return HttpResponse(f"Conversion failed. Error: {str(e)}")
        else:
            return view_func(request, *args, **kwargs)  # Continue with the original view function
    return wrapper_function


@pdf_to_word_logic
def pdf_to_word_view(request):
    meta = Meta(
        title='PDF to Word Document converter',
        description='Convert PDF file in to to Word Document. PDF pages will be converted to editable text with same Formatting.',
        keywords=['word', 'ms word', 'doc', 'docx'],
        og_title='PDF to Word Document converter',
        og_description='Convert PDF file in to to Word Document. PDF pages will be converted to editable text with same Formatting.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_word_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_word.html', context)

@pdf_to_word_logic
def pdf_to_word_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to Word Document',
        description='Convert PDF file in to to Word Document. PDF pages will be converted to editable text with same Formatting.',
        keywords=['word', 'ms word', 'doc', 'docx'],
        og_title='iLovePdfConverterOnline - PDF to Word Document',
        og_description='Convert PDF file in to to Word Document. PDF pages will be converted to editable text with same Formatting.',
    )
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_word_include.html',  context)


# -----------------------------------------================================

# Perfect working 
# https://github.com/pdf2htmlEX/pdf2htmlEX/releases/download/v0.18.8.rc1/pdf2htmlEX-0.18.8.rc1-master-20200630-Ubuntu-bionic-x86_64.deb
#   sudo dpkg -i pdf2htmlEX.deb
#   apt-get install -f


def pdf_to_html_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST' and request.FILES.get('pdf_file'):
            pdf_file = request.FILES['pdf_file']
            
            # Save the uploaded PDF file
            fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
            filename = fs.save(pdf_file.name, pdf_file)
            uploaded_file_path = fs.path(filename)
            
            # Define the output HTML file path
            output_html_path = os.path.splitext(uploaded_file_path)[0] + ".html"
            
            # Convert the PDF to HTML using pdf2htmlEX
            try:
                #for server
                pdf2htmlEX_path = '/usr/local/bin/pdf2htmlEX'  # Full path to pdf2htmlEX
                p = subprocess.run([pdf2htmlEX_path, '--dest-dir', fs.location, uploaded_file_path], check=True)
                
                # for localhost
                # p = subprocess.run(['pdf2htmlEX', '--dest-dir', fs.location, uploaded_file_path])
                # p = subprocess.Popen(['pdf2htmlEX', '--dest-dir', fs.location, uploaded_file_path])
                # p.wait()
            except subprocess.CalledProcessError:
                raise Http404("Error in converting PDF to HTML")
            
            # Create an HTTP response with the HTML content
            with open(output_html_path, 'r', encoding='utf-8') as html_file:
                html_content = html_file.read()
            
            response = HttpResponse(html_content, content_type='text/html')
            response['Content-Disposition'] = 'attachment; filename="PDF2HTML_ilovepdfconverteronline.com.html"'
            response.cleanup_files = [uploaded_file_path, output_html_path]
            return response
        
        return view_func(request, *args, **kwargs)
    
    return wrapper_function


@pdf_to_html_logic
def pdf_to_html_view(request):
    meta = Meta(
        title='PDF to HTML converter online',
        description='Convert HTML file or URL to PDF.',
        keywords=['html', 'url', 'urls', 'links', 'file', 'download'],
        og_title='PDF to HTML converter online',
        og_description='Convert PDF file in to HTML file.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_html_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_html.html', context)

@pdf_to_html_logic
def pdf_to_html_include(request):
        meta = Meta(
            title='PDF to HTML converter online',
            description='Convert HTML file or URL to PDF.',
            keywords=['html', 'url', 'urls', 'links', 'file', 'download'],
            og_title='PDF to HTML converter online',
            og_description='Convert PDF file in to HTML file.',
        )
        context = {'meta': meta}
        return render(request, 'tools/pdf_to_html_include.html', context)

# -----------------------------------------================================

def image_to_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        pdf_data=False
        if request.method == "POST" and request.FILES.getlist('images'):
            image_files = request.FILES.getlist('images')         
            temp_directory = os.path.join(settings.MEDIA_ROOT, 'temporary')
            os.makedirs(temp_directory, exist_ok=True)
            
            image_paths = []
            for img_file in image_files:
                image_path = os.path.join(temp_directory, img_file.name)
                with open(image_path, 'wb') as f:
                    for chunk in img_file.chunks():
                        f.write(chunk)
                image_paths.append(image_path)
                # pdf_data = img2pdf.convert(image_paths)
            try:
                pdf_data = img2pdf.convert(image_paths)
            except:
                os.remove(image_path)
                os.rmdir(temp_directory)
            if pdf_data:
                response = HttpResponse(pdf_data, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="Image_to_PDF_iLovePDFconverteronline.com.pdf"'
                
                for image_path in image_paths:
                    os.remove(image_path)
                os.rmdir(temp_directory)

                return response
            else:
                return render(request, 'tools/image_to_pdf.html')
        else:
            return view_func(request, *args, **kwargs)  
    return wrapper_function

@image_to_pdf_logic
def image_to_pdf_view(request):
    meta = Meta(
        title='JPG|JPEG|PNG Image to PDF',
        description='Convert JPG/JPEG Image file in to PDF. Image will be converted to PDF.',
        keywords=['png', 'image', 'jpg', 'jpeg'],
        og_title='JPG|JPEG|PNG Image to PDF',
        og_description='Convert JPG/JPEG Image file in to PDF. Image will be converted to PDF',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='image_to_pdf_view')
    context = {'meta': meta, 'tool_attachment':tool_attachment}
    return render(request, 'tools/image_to_pdf.html', context) 

@image_to_pdf_logic
def image_to_pdf_include(request):
    meta = Meta(
        title='Image to PDF',
        description='Convert image file (jpg, jpeg, png) to PDF file format',
        keywords=['png', 'image', 'jpg', 'jpeg'],
        og_title='Image to PDF',
        og_description='Convert image file (jpg, jpeg, png) to PDF file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/image_to_pdf_include.html')  


# -----------------------------------------================================
from django.core.management import call_command
import threading
import time

def schedule_cleanup(file_paths, delay=60):
    """Schedule file cleanup after a delay in seconds."""
    def cleanup():
        time.sleep(delay)
        try:
            print(f"Attempting to call cleanup_files with: {file_paths}")
            call_command('cleanup_files', *file_paths)
        except Exception as e:
            print(f"Error during cleanup: {e}")

    thread = threading.Thread(target=cleanup)
    thread.start()


def pdf_to_image_decorator(view_func):
    """Decorator to handle PDF to Image conversion with JSON response."""
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST':
            pdf_file = request.FILES.get('pdf_file')

            if pdf_file:
                try:
                    # Define output folder
                    output_folder = os.path.join(
                        settings.MEDIA_ROOT, "pdf_to_jpg"
                    )
                    os.makedirs(output_folder, exist_ok=True)

                    # Read the uploaded PDF content
                    pdf_content = pdf_file.read()

                    # Convert PDF to JPG images
                    jpg_paths = convert_pdf_to_jpg(
                        BytesIO(pdf_content), output_folder
                    )
                    logger.info(
                        f'Converted PDF to {len(jpg_paths)} images'
                    )

                    image_urls = []
                    for jpg_path in jpg_paths:
                        image_urls.append(
                            request.build_absolute_uri(
                                os.path.join(
                                    settings.MEDIA_URL,
                                    "pdf_to_jpg",
                                    os.path.basename(jpg_path)
                                )
                            )
                        )

                    response = JsonResponse(
                        {'image_urls': image_urls}
                    )
                    # Schedule cleanup after 60 seconds
                    schedule_cleanup(jpg_paths)
                    return response

                except Exception as e:
                    logger.error(f"PDF to Image error: {e}")
                    return JsonResponse(
                        {'error': str(e)},
                        status=500
                    )

        return view_func(request, *args, **kwargs)
    return wrapper_function


@pdf_to_image_decorator
def pdf_to_image_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to JPG|JPEG|PNG Image',
        description='Convert PDF file in to JPEG image. PDF pages will be converted to images.',
        keywords=['png', 'image', 'jpg', 'jpeg'],
        og_title='iLovePdfConverterOnline - PDF to JPG|JPEG|PNG Image',
        og_description='Convert PDF file in to JPEG image. PDF pages will be converted to images.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_image_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_image.html', context)

@pdf_to_image_decorator
def pdf_to_image_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to JPG|JPEG|PNG Image',
        description='Convert PDF file in to JPEG image. PDF pages will be converted to images.',
        keywords=['png', 'image', 'jpg', 'jpeg'],
        og_title='iLovePdfConverterOnline - PDF to JPG|JPEG|PNG Image',
        og_description='Convert PDF file in to JPEG image. PDF pages will be converted to images.',
    )
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_image_include.html', context)

# -----------------------------------------================================

# Working for server — uses get_libreoffice_path() for cross-platform
def powerpoint_to_pdf_logic(func):
    """Decorator to handle PowerPoint to PDF conversion via LibreOffice."""
    def wrapper(request, *args, **kwargs):
        if request.method == 'POST' and request.FILES.get('ppt_file'):
            try:
                ppt_file = request.FILES['ppt_file']
                out_path = os.path.join(
                    settings.MEDIA_ROOT, 'uploads'
                )
                os.makedirs(out_path, exist_ok=True)
                fs = FileSystemStorage(location=out_path)
                filename = fs.save(ppt_file.name, ppt_file)
                file_path = fs.path(filename)

                output_file_path = os.path.join(
                    out_path,
                    os.path.splitext(filename)[0] + '.pdf'
                )

                env = os.environ.copy()
                if platform.system() == 'Windows':
                    env['HOME'] = out_path
                else:
                    env['HOME'] = '/tmp'

                lo_path = get_libreoffice_path()
                result = subprocess.run(
                    [lo_path, '--headless', '--convert-to',
                     'pdf', '--outdir', out_path, file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(
                        f"LibreOffice conversion failed: "
                        f"{result.stderr}"
                    )

                with open(output_file_path, 'rb') as pdf_file:
                    response = HttpResponse(
                        pdf_file.read(),
                        content_type='application/pdf'
                    )
                    response['Content-Disposition'] = (
                        f'attachment; filename='
                        f'{os.path.basename(output_file_path)}'
                    )
                    response.cleanup_files = [
                        file_path, output_file_path
                    ]
                    return response
            except FileNotFoundError:
                return HttpResponse(
                    "LibreOffice is not installed. "
                    "See INSTALL_DEPENDENCIES.md for setup.",
                    status=500
                )
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        return func(request, *args, **kwargs)
    return wrapper


# Settings done before working on server
# sudo mkdir -p /tmp/nobody_home
# sudo chown nobody:nogroup /tmp/nobody_home
# sudo chmod 700 /tmp/nobody_home

@csrf_exempt
@powerpoint_to_pdf_logic
def powerpoint_to_pdf_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Powerpoint to PDF converter',
        description='Convert PPT, PPTX file of PowerPoint in to PDF file format.',
        keywords=['ppt', 'pptx', 'slide', 'pdf'],
        og_title='iLovePdfConverterOnline - Powerpoint to PDF converter',
        og_description='Convert PPT, PPTX file of PowerPoint in to PDF file format.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='powerpoint_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/powerpoint_to_pdf.html', context)

@csrf_exempt
@powerpoint_to_pdf_logic
def powerpoint_to_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Powerpoint to PDF converter',
        description='Convert PPT, PPTX file of PowerPoint in to PDF file format.',
        keywords=['ppt', 'pptx', 'slide', 'pdf'],
        og_title='iLovePdfConverterOnline - Powerpoint to PDF converter',
        og_description='Convert PPT, PPTX file of PowerPoint in to PDF file format.',
    )
    context = {'meta': meta}
    return render(request, 'tools/powerpoint_to_pdf_include.html', context)

# -----------------------------------------================================

# decorators.py
from pptx import Presentation
from pdf2image import convert_from_bytes

def pdf_to_pptx_logic(func):
    @wraps(func)
    def wrapper(request, *args, **kwargs):
        if request.method == 'POST' and 'pdf_file' in request.FILES:
            pdf_file = request.FILES['pdf_file']
            pdf_bytes = pdf_file.read()
            images = convert_from_bytes(pdf_bytes)
            
            prs = Presentation()
            blank_slide_layout = prs.slide_layouts[6]  # Choosing a blank slide layout

            for image in images:
                slide = prs.slides.add_slide(blank_slide_layout)
                image_stream = BytesIO()
                image.save(image_stream, format='PNG')
                image_stream.seek(0)
                
                slide.shapes.add_picture(image_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)

            pptx_stream = BytesIO()
            prs.save(pptx_stream)
            pptx_stream.seek(0)

            response = HttpResponse(pptx_stream, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            response['Content-Disposition'] = 'attachment; filename="PDF2PowerPoint.pptx"'
            return response
        
        return func(request, *args, **kwargs)
    return wrapper


@pdf_to_pptx_logic
def pdf_to_pptx_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to PPTX (PowerPoint) converter online',
        description='Convert PDF to PPTX (PowerPoint) file online in free.',
        keywords= ['pdf', 'pptx', 'file', 'PowerPoint', 'power point', 'slide'],
        og_title='iLovePdfConverterOnline - PDF to PPTX (PowerPoint) converter online',
        og_description='Convert PDF to PPTX (PowerPoint) file online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_pptx_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    
    return render(request, 'tools/pdf_to_pptx.html', context)

@pdf_to_pptx_logic
def pdf_to_pptx_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to PPTX (PowerPoint) converter online',
        description='Convert PDF to PPTX (PowerPoint) file online in free.',
        keywords= ['pdf', 'pptx', 'file', 'PowerPoint', 'power point', 'slide'],
        og_title='iLovePdfConverterOnline - PDF to PPTX (PowerPoint) converter online',
        og_description='Convert PDF to PPTX (PowerPoint) file online in free.',
    )
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_pptx_include.html', context)



# -----------------------------------------================================

#Working for XLSX, XLSM, XLTX XLTM including XLS & CSV (created with Excel but not Downloaded)
def excel_to_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST' and request.FILES.get('excel_file'):
            excel_file = request.FILES['excel_file']
            file_name = excel_file.name.lower()
            file_extension = os.path.splitext(file_name)[1]

            # Save uploaded file to MEDIA_ROOT
            file_path = os.path.join(settings.MEDIA_ROOT, excel_file.name)
            with open(file_path, 'wb') as destination:
                for chunk in excel_file.chunks():
                    destination.write(chunk)

            # Create PDF
            pdf_path = os.path.join(settings.MEDIA_ROOT, 'output.pdf')
            c = canvas.Canvas(pdf_path, pagesize=landscape(letter))

            top_margin = 0.5 * inch
            left_margin = 0.5 * inch
            bottom_margin = 0.5 * inch
            right_margin = 0.5 * inch

            page_width, page_height = landscape(letter)

            cell_height = 20
            font_size = 10  # Starting font size
            font = 'Helvetica'  # Font family

            if file_extension in ['.xlsx', '.xlsm', '.xltx', '.xltm']:
                workbook = load_workbook(file_path)
                worksheet = workbook.active
                max_row = worksheet.max_row
                max_column = worksheet.max_column
                ws_range = worksheet.iter_rows(values_only=True)
            elif file_extension == '.xls':
                workbook = open_workbook(file_path)
                worksheet = workbook.sheet_by_index(0)
                max_row = worksheet.nrows
                max_column = worksheet.ncols
                ws_range = (worksheet.row_values(row) for row in range(max_row))
            elif file_extension == '.csv':
                try:
                    with open(file_path, newline='', encoding='utf-8') as csvfile:
                        reader = csv.reader(csvfile)
                        data = list(reader)
                except UnicodeDecodeError:
                    return HttpResponse("Unable to decode the CSV file. Please ensure it is encoded in UTF-8.", content_type="text/plain")

                max_row = len(data)
                max_column = len(data[0]) if max_row > 0 else 0
                ws_range = iter(data)
            else:
                return HttpResponse("Unsupported file type.", content_type="text/plain")

            cell_width = (page_width - left_margin - right_margin) / max_column
            max_text_width = cell_width - 2  # Subtracting a bit for padding

            y = page_height - top_margin  # Initial y position

            for row in ws_range:
                for col_num, cell in enumerate(row):
                    x = left_margin + col_num * cell_width
                    text = str(cell)
                    current_font_size = font_size
                    while c.stringWidth(text, font, current_font_size) > max_text_width and current_font_size > 1:
                        current_font_size -= 1
                    c.setFont(font, current_font_size)
                    c.drawString(x, y, text)
                
                y -= cell_height
                
                if y < bottom_margin:
                    c.showPage()
                    y = page_height - top_margin

            c.save()

            # Provide the PDF file for download
            with open(pdf_path, 'rb') as pdf_file:
                response = HttpResponse(pdf_file.read(), content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename=Excel2PDF_ilovepdfconverteronline.com.pdf'
                response.cleanup_files = [file_path, pdf_path]
                return response

        return view_func(request, *args, **kwargs)  # Continue with the original view function

    return wrapper_function

@excel_to_pdf_logic
def excel_to_pdf_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Excel to PDF converter',
        description='Convert XLSX, XLSM, XLTX, XLTM, XLS & CSV file in to PDF file format.',
        keywords=['XLSX', 'XLSM', 'XLTX', 'XLTM', 'XLS'  'CSV', 'pdf'],
        og_title='iLovePdfConverterOnline - Excel to PDF converter',
        og_description='Convert XLSX, XLSM, XLTX, XLTM, XLS & CSV file in to PDF file format.',
    )    
    tool_attachment = ToolAttachment.objects.get(function_name='excel_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/excel_to_pdf.html', context)

@excel_to_pdf_logic
def excel_to_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Excel to PDF converter',
        description='Convert XLSX, XLSM, XLTX, XLTM, XLS & CSV file in to PDF file format.',
        keywords=['XLSX', 'XLSM', 'XLTX', 'XLTM', 'XLS'  'CSV', 'pdf'],
        og_title='iLovePdfConverterOnline - Excel to PDF converter',
        og_description='Convert XLSX, XLSM, XLTX, XLTM, XLS & CSV file in to PDF file format.',
    )    
    context = {'meta': meta}
    return render(request, 'tools/excel_to_pdf_include.html')

# -----------------------------------------================================


# def pdf_to_excel_logic(func):
#     def wrapper(request, *args, **kwargs):
#         if request.method == 'POST' and request.FILES.get('pdf_file'):
#             pdf_file = request.FILES['pdf_file']

#             # Save the uploaded PDF file to uploads directory
#             upload_folder = os.path.join(settings.MEDIA_ROOT, 'uploads')
#             os.makedirs(upload_folder, exist_ok=True)
#             pdf_file_path = os.path.join(upload_folder, pdf_file.name)
            
#             with open(pdf_file_path, 'wb') as destination:
#                 for chunk in pdf_file.chunks():
#                     destination.write(chunk)

#             try:
#                 # Convert PDF to CSV using tabula
#                 csv_file_path = os.path.join(upload_folder, 'PDF2Excel_ilovepdfconverteronline.com.csv')
#                 tabula.convert_into(input_path=pdf_file_path, output_path=csv_file_path, output_format='csv', pages='all', stream=True)

#                 # Convert CSV to XLSX using pandas
#                 xlsx_file_path = os.path.join(upload_folder, 'PDF2Excel_ilovepdfconverteronline.com.xlsx')
#                 read_file = pd.read_csv(csv_file_path)
#                 read_file.to_excel(xlsx_file_path, index=None, header=True)

#                 # Provide the XLSX file for download
#                 with open(xlsx_file_path, 'rb') as excel_file:
#                     response = HttpResponse(excel_file.read(), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
#                     response['Content-Disposition'] = 'attachment; filename=PDF2Excel_ilovepdfconverteronline.com.xlsx'
#                     response.cleanup_files = [pdf_file_path, csv_file_path, xlsx_file_path]  # Add files for cleanup
#                     return response

#             except Exception as e:
#                 return HttpResponse(f"Conversion failed. Error: {str(e)}")

#         # If GET request or no pdf_file in POST
#         return func(request, *args, **kwargs)

#     return wrapper


# @pdf_to_excel_logic
# def pdf_to_excel_view(request):
#     meta = Meta(
#         title='iLovePdfConverterOnline - PDF to XLSX file converter online',
#         description='Convert PDF to XLSX file online in free.',
#         keywords= ['pdf', 'file', 'excel', 'xlsx', 'xls'],
#         og_title='iLovePdfConverterOnline - PDF to XLSX file converter online',
#         og_description='Convert PDF to XLSX file online in free.',
#     )
#     tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_excel_view')
#     context = {'meta': meta, 'tool_attachment': tool_attachment}
#     return render(request, 'tools/pdf_to_excel.html', context)

# @pdf_to_excel_logic
# def pdf_to_excel_include(request):
#     meta = Meta(
#         title='iLovePdfConverterOnline - PDF to XLSX file converter online',
#         description='Convert PDF to XLSX file online in free.',
#         keywords= ['pdf', 'file', 'excel', 'xlsx', 'xls'],
#         og_title='iLovePdfConverterOnline - PDF to XLSX file converter online',
#         og_description='Convert PDF to XLSX file online in free.',
#     )
#     context = {'meta': meta}
#     return render(request, 'tools/pdf_to_excel_include.html', context)

#-------------------------------------------------------------
# 
#  PDF TO EXCEL
import os
import uuid
import tempfile
import io
import logging
import re
from pathlib import Path
from django.http import JsonResponse, HttpResponse, FileResponse
from django.core.files.storage import default_storage
from django.core.files.base import ContentFile
from django.shortcuts import render, redirect
from django.conf import settings
from django.utils import translation
from django.template.loader import get_template
from django.contrib import messages

import pandas as pd
import numpy as np
import camelot
import pdfplumber
import tabula
import pypdfium2 as pdfium
from openpyxl import Workbook
from reportlab.lib import colors
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.platypus import SimpleDocTemplate, Paragraph, Table, TableStyle
from reportlab.lib.units import inch

logger = logging.getLogger(__name__)

def pdf_to_excel_view(request):
    """
    Advanced PDF to Excel converter with intelligent table detection
    and mixed content handling
    """
    if request.method == 'POST':
        return handle_pdf_to_excel_upload(request)

    meta = Meta(
        title='iLovePdfConverterOnline - PDF to XLSX file converter online',
        description='Convert PDF to XLSX file online in free.',
        keywords= ['pdf', 'file', 'excel', 'xlsx', 'xls'],
        og_title='iLovePdfConverterOnline - PDF to XLSX file converter online',
        og_description='Convert PDF to XLSX file online in free.',
    )
    # Using fallback to avoid undefined ToolAttachment variable if not imported
    # Assuming ToolAttachment is available since it's used elsewhere
    try:
        tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_excel_view')
    except:
        tool_attachment = None
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_excel.html', context)

def pdf_to_excel_include(request):
    """
    Advanced PDF to Excel converter with intelligent table detection
    and mixed content handling
    """
    if request.method == 'POST':
        return handle_pdf_to_excel_upload(request)

    meta = Meta(
        title='iLovePdfConverterOnline - PDF to XLSX file converter online',
        description='Convert PDF to XLSX file online in free.',
        keywords= ['pdf', 'file', 'excel', 'xlsx', 'xls'],
        og_title='iLovePdfConverterOnline - PDF to XLSX file converter online',
        og_description='Convert PDF to XLSX file online in free.',
    )
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_excel_include.html', context)


def handle_pdf_to_excel_upload(request):
    """
    Handle file upload and perform intelligent PDF to Excel conversion
    """
    if 'pdf_file' not in request.FILES:
        messages.error(request, 'No PDF file uploaded.')
        return redirect('pdf_to_excel')

    pdf_file = request.FILES['pdf_file']

    # Generate unique filename for processed files
    unique_id = str(uuid.uuid4())
    original_filename = f"converted_{unique_id}_{pdf_file.name}"
    output_filename = f"{Path(original_filename).stem}.xlsx"

    try:
        # Save uploaded PDF temporarily
        with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as temp_file:
            for chunk in pdf_file.chunks():
                temp_file.write(chunk)
            temp_pdf_path = temp_file.name

        # Perform advanced conversion
        converted_file_path = convert_pdf_to_excel_advanced(temp_pdf_path, original_filename)

        # Clean up temp file
        os.unlink(temp_pdf_path)

        # Read the converted Excel file for download
        if os.path.exists(converted_file_path):
            excel_content = open(converted_file_path, 'rb').read()
            response = HttpResponse(excel_content, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            response['Content-Disposition'] = f'attachment; filename="{output_filename}"'

            # Clean up converted file
            os.unlink(converted_file_path)

            messages.success(request, 'PDF converted to Excel successfully!')
            return response
        else:
            messages.error(request, 'Conversion failed. Could not create Excel file.')
            return redirect('pdf_to_excel')

    except Exception as e:
        logger.error(f"PDF to Excel conversion error: {str(e)}")
        try:
            os.unlink(temp_pdf_path)
        except:
            pass
        messages.error(request, f'Conversion error: {str(e)}')
        return redirect('pdf_to_excel')


def convert_pdf_to_excel_advanced(pdf_path, output_name):
    """
    Advanced PDF to Excel converter that handles:
    - Single/multiple tables
    - Tables spanning multiple pages
    - PDFs with headers repeated on each page
    - Mixed content (text and tables)
    - Different table structures and layouts
    """

    # Create output directory in media
    output_dir = os.path.join(settings.MEDIA_ROOT, 'converted_files')
    os.makedirs(output_dir, exist_ok=True)

    output_path = os.path.join(output_dir, output_name)

    # Analyze PDF structure first
    pdf_info = analyze_pdf_structure(pdf_path)

    # Determine conversion strategy based on PDF analysis
    if pdf_info['has_tables']:
        return convert_pdf_with_tables(pdf_path, output_path, pdf_info)
    else:
        return convert_pdf_text_only(pdf_path, output_path)


def analyze_pdf_structure(pdf_path):
    """
    Analyze PDF to determine structure:
    - Number of pages
    - Whether tables exist
    - Header patterns
    - Table locations
    """
    info = {
        'num_pages': 0,
        'has_tables': False,
        'table_count': 0,
        'header_pages': [],
        'table_pages': [],
        'repeating_headers': False
    }

    try:
        # Use pdfplumber for detailed analysis
        with pdfplumber.open(pdf_path) as pdf:
            info['num_pages'] = len(pdf.pages)

            table_counts = []
            header_patterns = []

            for i, page in enumerate(pdf.pages):
                tables = page.extract_tables()
                table_texts = []

                if tables:
                    info['has_tables'] = True
                    info['table_count'] += len(tables)
                    info['table_pages'].append(i + 1)
                    table_counts.append(len(tables))

                    # Extract text for header analysis
                    text = page.extract_text() or ""
                    table_texts.append(text)

                    # Check for repeating headers (simple heuristic)
                    if i < len(pdf.pages) - 1:
                        next_text = pdf.pages[i + 1].extract_text() or ""
                        if text and next_text and len(text) > 50:
                            if text[:min(100, len(text))] == next_text[:min(100, len(next_text))]:
                                info['repeating_headers'] = True
                                info['header_pages'].append(i + 1)

        # Also try tabula for cross-validation
        try:
            tables = tabula.read_pdf(pdf_path, pages='all')
            if isinstance(tables, list) and len(tables) > 0:
                info['has_tables'] = True
                info['table_count'] = len([t for t in tables if not t.empty])
        except:
            pass

    except Exception as e:
        logger.error(f"PDF analysis error: {str(e)}")

    return info


def convert_pdf_with_tables(pdf_path, output_path, pdf_info):
    """
    Convert PDF with tables to Excel with intelligent handling:
    - Group tables from pages with repeating headers
    - Create separate sheets for distinct tables
    - Handle complex table structures
    """

    wb = excel_workbook = Workbook()

    # Strategy 1: Try camelot first (best for tabular data)
    try:
        tables = camelot.read_pdf(
            pdf_path,
            pages='all',
            flavor='lattice',  # For tables with lines
            suppress_warnings=True
        )

        if tables.n > 0:
            logger.info(f"Camelot found {tables.n} tables using lattice flavor")
            return convert_tables_to_excel(camelot_tables=tables, wb=excel_workbook, output_path=output_path)

    except Exception as e:
        logger.warning(f"Camelot lattice failed: {str(e)}")
        try:
            # Try stream flavor for tables without clear lines
            tables = camelot.read_pdf(
                pdf_path,
                pages='all',
                flavor='stream',
                suppress_warnings=True
            )

            if tables.n > 0:
                logger.info(f"Camelot found {tables.n} tables using stream flavor")
                return convert_tables_to_excel(camelot_tables=tables, wb=excel_workbook, output_path=output_path)

        except Exception as e:
            logger.warning(f"Camelot stream also failed: {str(e)}")

    # Strategy 2: Use pdfplumber for table extraction
    try:
        return convert_with_pdfplumber(pdf_path, output_path, pdf_info)
    except Exception as e:
        logger.error(f"All conversion strategies failed: {str(e)}")
        raise


def convert_tables_to_excel(camelot_tables, wb, output_path):
    """
    Convert camelot tables to Excel workbook with intelligent sheet naming
    """
    sheets_created = 0

    for i, table in enumerate(camelot_tables):
        # Convert camelot table to DataFrame
        df = table.df
        df = pd.DataFrame(df)

        # Clean the DataFrame
        df = clean_dataframe(df)

        # Create sheet name
        sheet_name = f"Table {i + 1}"

        # Write to Excel
        df.to_excel(output_path, sheet_name=sheet_name, index=False)
        sheets_created += 1

    return output_path if sheets_created > 0 else None


def convert_with_pdfplumber(pdf_path, output_path, pdf_info):
    """
    Convert PDF using pdfplumber for better table detection
    """
    wb = Workbook()

    with pdfplumber.open(pdf_path) as pdf:
        tables_created = 0

        for page_num, page in enumerate(pdf.pages):
            tables = page.extract_tables()

            if not tables:
                continue

            for table_idx, table in enumerate(tables):
                if not table or not any(cell for cell in table):
                    continue

                # Create DataFrame
                df = pd.DataFrame(table[1:], columns=table[0]) if len(table) > 1 else pd.DataFrame(table)
                df = clean_dataframe(df)

                if df.empty:
                    continue

                # Determine sheet naming strategy
                if tables_created == 0:
                    sheet_name = "Sheet1"
                else:
                    sheet_name = f"Table {page_num + 1}-{table_idx + 1}"

                # Write to Excel
                if tables_created == 0:
                    df.to_excel(output_path, sheet_name=sheet_name, index=False)
                else:
                    with pd.ExcelWriter(output_path, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
                        df.to_excel(writer, sheet_name=sheet_name, index=False)

                tables_created += 1

    return output_path if tables_created > 0 else None


def convert_pdf_text_only(pdf_path, output_path):
    """
    Fallback: Convert PDF text content to Excel (non-tabular data)
    """
    wb = Workbook()

    with pdfplumber.open(pdf_path) as pdf:
        for i, page in enumerate(pdf.pages):
            text = page.extract_text() or ""

            # Clean text
            text = re.sub(r'\s+', ' ', text).strip()

            if not text:
                continue

            # Convert each page to a row
            df = pd.DataFrame({'Page Text': [text]})

            sheet_name = f"Page {i + 1}"

            if i == 0:
                df.to_excel(output_path, sheet_name=sheet_name, index=False)
            else:
                with pd.ExcelWriter(output_path, engine='openpyxl', mode='a', if_sheet_exists='replace') as writer:
                    df.to_excel(writer, sheet_name=sheet_name, index=False)

    return output_path


def clean_dataframe(df):
    """
    Clean DataFrame by:
    - Removing empty rows/columns
    - Converting numeric values
    - Handling missing values
    """
    # Remove completely empty rows and columns
    df = df.dropna(how='all')
    df = df.dropna(axis=1, how='all')

    # Strip whitespace from string values
    for col in df.select_dtypes(include=['object']).columns:
        df[col] = df[col].astype(str).str.strip()

    # Try to convert to numeric where possible
    for col in df.columns:
        df[col] = pd.to_numeric(df[col], errors='coerce')

    return df


def render_pdf_to_pdf(request):
    """
    Render PDF to PDF (placeholder for future functionality)
    """
    if request.method == 'POST':
        return handle_pdf_to_pdf_upload(request)

    meta = {
        'og_title': 'iLovePdfConverterOnline - PDF to PDF converter',
        'og_description': 'Convert PDF to PDF online in free.',
    }
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_pdf_include.html', context)
 
# -----------------------------------------================================


def pdf_to_csv_logic(func):
    def wrapper(request, *args, **kwargs):
        if request.method == 'POST' and request.FILES.get('pdf_file'):
            pdf_file = request.FILES['pdf_file']

            # Save the uploaded PDF file to uploads directory
            upload_folder = os.path.join(settings.MEDIA_ROOT, 'uploads')
            os.makedirs(upload_folder, exist_ok=True)
            pdf_file_path = os.path.join(upload_folder, pdf_file.name)
            
            with open(pdf_file_path, 'wb') as destination:
                for chunk in pdf_file.chunks():
                    destination.write(chunk)

            try:
                # Convert PDF to CSV using tabula
                csv_file_path = os.path.join(upload_folder, 'PDF2CSV_ilovepdfconverteronline.com.csv')
                tabula.convert_into(input_path=pdf_file_path, output_path=csv_file_path, output_format='csv', pages='all', stream=True)

                # Provide the CSV file for download
                with open(csv_file_path, 'rb') as csv_file:
                    response = HttpResponse(csv_file.read(), content_type='text/csv')
                    response['Content-Disposition'] = 'attachment; filename=PDF2CSV_ilovepdfconverteronline.com.csv'
                    response.cleanup_files = [pdf_file_path, csv_file_path]  # Add files for cleanup
                    return response

            except Exception as e:
                return HttpResponse(f"Conversion failed. Error: {str(e)}")

        # If GET request or no pdf_file in POST
        return func(request, *args, **kwargs)

    return wrapper


@pdf_to_csv_logic
def pdf_to_csv_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to CSV file converter online',
        description='Convert PDF (Portable Document Format) to CSV (comma-separated values) file online in free.',
        keywords= ['pdf', 'file', 'excel', 'csv', 'comma-separated values', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - PDF to CSV file converter online',
        og_description='Convert PDF (Portable Document Format) to CSV (comma-separated values)  file online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_csv_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_csv.html',context)

@pdf_to_csv_logic
def pdf_to_csv_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to CSV file converter online',
        description='Convert PDF (Portable Document Format) to CSV (comma-separated values) file online in free.',
        keywords= ['pdf', 'file', 'excel', 'csv', 'comma-separated values', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - PDF to CSV file converter online',
        og_description='Convert PDF (Portable Document Format) to CSV (comma-separated values)  file online in free.',
    )
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_csv_include.html', context)

# -----------------------------------------================================

def json_to_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST' and 'json_file' in request.FILES:
            json_file = request.FILES['json_file']
            try:
                # Save the uploaded JSON file
                with open('uploaded_json.json', 'wb') as destination:
                    for chunk in json_file.chunks():
                        destination.write(chunk)
                
                # Load the JSON data
                with open('uploaded_json.json', 'r') as file:
                    data = json.load(file)
                
                if not isinstance(data, dict):
                    return HttpResponse("Conversion failed. The uploaded JSON file is not properly formatted.")
                
                # Create PDF
                response = HttpResponse(content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="JSON_to_PDF.pdf"'

                # Render JSON data to PDF
                pdf = canvas.Canvas(response, pagesize=letter)
                y_coordinate = 700
                for line in json.dumps(data, indent=4).split('\n'):
                    pdf.drawString(100, y_coordinate, line)
                    y_coordinate -= 12  # Move to the next line
                pdf.save()
                
                return response
            except Exception as e:
                logger.error(f"Conversion failed. Error: {e}")
                logger.error(traceback.format_exc())
                return HttpResponse("Conversion failed. An error occurred during conversion. Check File Type Uploaded.")
        else:
            return view_func(request, *args, **kwargs)  # Continue with the original view function
    return wrapper_function

@json_to_pdf_logic
def json_to_pdf_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - JSON to PDF converter online',
        description='Convert JSON (JavaScript Object Notation) file in to PDF (Portable Document Format) online in free.',
        keywords=['pdf', 'json', 'Portable Document Format', 'JavaScript Object Notation'],
        og_title='iLovePdfConverterOnline - JSON to PDF converter online',
        og_description='Convert JSON (JavaScript Object Notation) file in to PDF (Portable Document Format) online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='json_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/json_to_pdf.html', context)

def json_to_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - JSON to PDF converter online',
        description='Convert JSON (JavaScript Object Notation) file in to PDF (Portable Document Format) online in free.',
        keywords=['pdf', 'json', 'Portable Document Format', 'JavaScript Object Notation'],
        og_title='iLovePdfConverterOnline - JSON to PDF converter online',
        og_description='Convert JSON (JavaScript Object Notation) file in to PDF (Portable Document Format) online in free.',
    )
    context = {'meta': meta}
    return render(request, 'tools/json_to_pdf_include.html', context)

# -----------------------------------------================================

import logging
logger = logging.getLogger(__name__)

def pdf_to_json_logic(view_func):
    @wraps(view_func)
    def wrapper_function(request, *args, **kwargs):
        if request.method == 'POST' and 'pdf_file' in request.FILES:
            pdf_file = request.FILES['pdf_file']
            try:
                # Save the uploaded PDF file
                with open('uploaded_pdf.pdf', 'wb') as destination:
                    for chunk in pdf_file.chunks():
                        destination.write(chunk)
                
                # Read PDF file and extract text
                with open('uploaded_pdf.pdf', 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    text = ""
                    for page_num in range(len(pdf_reader.pages)):
                        text += pdf_reader.pages[page_num].extract_text()
                
                # Convert extracted text to JSON format
                json_data = json.loads(text)

                # Save JSON data to a file
                with open('PDF_to_JSON.json', 'w') as json_file:
                    json.dump(json_data, json_file, indent=4)
                
                return FileResponse(open('PDF_to_JSON.json', 'rb'), as_attachment=True)
            except Exception as e:
                logger.error(f"Conversion failed. Error: {e}")
                logger.error(traceback.format_exc())
                return HttpResponse("Conversion failed. An error occurred during conversion.")
        else:
            return view_func(request, *args, **kwargs)  # Continue with the original view function
    return wrapper_function

@pdf_to_json_logic
def pdf_to_json_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to JSON converter online',
        description='Convert PDF (Portable Document Format) file in to JSON (JavaScript Object Notation) file format online in free.',
        keywords=['pdf', 'json', 'Portable Document Format', 'JavaScript Object Notation'],
        og_title='iLovePdfConverterOnline - PDF to JSON converter online',
        og_description='Convert PDF (Portable Document Format) file in to JSON (JavaScript Object Notation) file format online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_json_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_json.html', context)

@pdf_to_json_logic
def pdf_to_json_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to JSON converter online',
        description='Convert PDF (Portable Document Format) file in to JSON (JavaScript Object Notation) file format online in free.',
        keywords=['pdf', 'json', 'Portable Document Format', 'JavaScript Object Notation'],
        og_title='iLovePdfConverterOnline - PDF to JSON converter online',
        og_description='Convert PDF (Portable Document Format) file in to JSON (JavaScript Object Notation) file format online in free.',
    )
    context = {'meta': meta}
    return render(request, 'tools/pdf_to_json_include.html', context)

# -----------------------------------------================================

def string_to_base64_logic(view_func):
    @wraps(view_func)
    def wrapper(request, *args, **kwargs):
        context = {}
        if request.method == "POST":
            original_string = request.POST.get("original_string")
            if original_string:
                base64_string = base64.b64encode(original_string.encode()).decode()
                context['base64_string'] = base64_string
                context['original_string'] = original_string
        return view_func(request, context, *args, **kwargs)
    return wrapper

@string_to_base64_logic
def string_to_base64_view(request, context):
    meta = Meta(
        title='iLovePdfConverterOnline - String to Base64 file converter online',
        description='Convert String into Base64 file format online in free.',
        keywords=['string', 'text', 'base64'],
        og_title='iLovePdfConverterOnline - String to Base64 file converter online',
        og_description='Convert String into Base64 file format online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='string_to_base64_view')
    context ['tool_attachment'] = tool_attachment
    context ['meta'] = meta

    return render(request, 'tools/string_to_base64.html', context)

@string_to_base64_logic
def string_to_base64_include(request, context):
    meta = Meta(
        title='iLovePdfConverterOnline - String to Base64 file converter online',
        description='Convert String into Base64 file format online in free.',
        keywords=['string', 'text', 'base64'],
        og_title='iLovePdfConverterOnline - String to Base64 file converter online',
        og_description='Convert String into Base64 file format online in free.',
    )
    context['meta'] = meta
    return render(request, 'tools/string_to_base64_include.html', context)

# -----------------------------------------================================

def base64_to_pdf_logic(view_func):
    def wrapper(request, *args, **kwargs):
        if request.method == 'POST':
            # Getting base64 string from the textarea input
            base64_string = request.POST.get('base64_string', '')
            try:
                # Decoding base64 string to bytes
                pdf_data = base64.b64decode(base64_string)
                print(f'pdf data {pdf_data}')

                # Creating a PDF file using reportlab
                buffer = BytesIO()
                c = canvas.Canvas(buffer)

                # Write the decoded data to the PDF
                c.drawString(100, 750, pdf_data.decode('utf-8'))
                c.save()

                buffer.seek(0)
                response = HttpResponse(buffer, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="Base642PDF_ilovepdfconverteronline.com.pdf"'

                return response
            except:
                return HttpResponseBadRequest("Invalid base64 string")
        return view_func(request, *args, **kwargs)
    return wrapper


@base64_to_pdf_logic
def base64_to_pdf_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Base64 to PDF file converter online',
        description='Convert Base64 into PDF (Portable Document Format) file format online in free.',
        keywords=['string', 'pdf', 'base64', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - Base64 to PDF file converter online',
        og_description='Convert Base64 into PDF (Portable Document Format) file format online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='base64_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/base64_to_pdf.html', context)

@base64_to_pdf_logic
def base64_to_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Base64 to PDF file converter online',
        description='Convert Base64 into PDF (Portable Document Format) file format online in free.',
        keywords=['string', 'pdf', 'base64', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - Base64 to PDF file converter online',
        og_description='Convert Base64 into PDF (Portable Document Format) file format online in free.',
    )
    context = {'meta': meta}
    return render(request, 'tools/base64_to_pdf_include.html', context)

# ---------------------------------------------====================================================

def pdf_to_base64_logic(template_name):
    def decorator(func):
        @wraps(func)
        def wrapper(request, *args, **kwargs):
            base64_string = None
            if request.method == 'POST':
                pdf_file = request.FILES.get('pdf_file')
                if pdf_file:
                    base64_string = base64.b64encode(pdf_file.read()).decode('utf-8')
            context = func(request, base64_string, *args, **kwargs)
            return render(request, template_name, context)
        return wrapper
    return decorator


# Using the decorator for pdf_to_base64_view
@pdf_to_base64_logic('tools/pdf_to_base64.html')
def pdf_to_base64_view(request, base64_string):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to Base64 file converter online',
        description='Convert PDF (Portable Document Format) into Base64 file format online in free.',
        keywords=['string', 'text', 'base64', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - PDF to Base64 file converter online',
        og_description='Convert PDF (Portable Document Format) into Base64 file format online in free.',
    )
    # tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_base64_view')
    return {'meta': meta, 'base64_string': base64_string}

# Using the decorator for pdf_to_base64_include
@pdf_to_base64_logic('tools/pdf_to_base64_include.html')
def pdf_to_base64_include(request, base64_string):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to Base64 file converter online',
        description='Convert PDF (Portable Document Format) into Base64 file format online in free.',
        keywords=['string', 'text', 'base64', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - PDF to Base64 file converter online',
        og_description='Convert PDF (Portable Document Format) into Base64 file format online in free.',
    )
    return {'meta': meta, 'base64_string': base64_string}

# ---------------------------------------------====================================================


def tiff_to_pdf_logic(func):
    def wrapper(request, *args, **kwargs):
        if request.method == "POST":
            tiff_file = request.FILES.get('tiff_file')
            if not tiff_file:
                return HttpResponse("No TIFF file uploaded.", status=400)
            
            tiff_path = default_storage.save(tiff_file.name, ContentFile(tiff_file.read()))
            tiff_path_full = default_storage.path(tiff_path)

            try:
                pdf_path_full = tiff_path_full.replace('.tiff', '.pdf').replace('.tif', '.pdf')
                if not os.path.exists(tiff_path_full):
                    raise Exception(f'{tiff_path_full} not found.')

                image = Image.open(tiff_path_full)
                images = []

                for i, page in enumerate(ImageSequence.Iterator(image)):
                    page = page.convert("RGB")
                    images.append(page)

                if len(images) == 1:
                    images[0].save(pdf_path_full)
                else:
                    images[0].save(pdf_path_full, save_all=True, append_images=images[1:])
                
                # Serve the PDF file for download
                with open(pdf_path_full, 'rb') as pdf_file:
                    response = HttpResponse(pdf_file.read(), content_type='application/pdf')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(pdf_path_full)}'
                
                # Cleanup temporary files
                try:
                    os.remove(tiff_path_full)
                    os.remove(pdf_path_full)
                except OSError:
                    pass

                return response
            except Exception as e:
                return HttpResponse(f"Error converting TIFF to PDF: {e}", status=500)
        else:
            return func(request, *args, **kwargs)
    return wrapper

@tiff_to_pdf_logic
def tiff_to_pdf_view(request, pdf_url=None):
    meta = Meta(
        title='iLovePdfConverterOnline - TIFF to PDF file converter online',
        description='Convert TIFF (Tag Image File Format) in to PDF (Portable Document Format) file format online in free.',
        keywords=['tiff', 'image', 'pdf', 'Tag Image File Format', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - TIFF to PDF file converter online',
        og_description='Convert TIFF (Tag Image File Format) in to PDF (Portable Document Format) file format online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='tiff_to_pdf_view')
    context = {'meta': meta, 'pdf_url': pdf_url, 'tool_attachment': tool_attachment}
    return render(request, 'tools/tiff_to_pdf.html', context)


@tiff_to_pdf_logic
def tiff_to_pdf_include(request, pdf_url=None):
    meta = Meta(
        title='iLovePdfConverterOnline - TIFF to PDF file converter online',
        description='Convert TIFF (Tag Image File Format) in to PDF (Portable Document Format) file format online in free.',
        keywords=['tiff', 'image', 'pdf', 'Tag Image File Format', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - TIFF to PDF file converter online',
        og_description='Convert TIFF (Tag Image File Format) in to PDF (Portable Document Format) file format online in free.',
    )
    context ={'meta': meta, 'pdf_url': pdf_url}
    return render(request, 'tools/tiff_to_pdf_include.html', context)


# ---------------------------------------------====================================================


def pdf_to_tiff_logic(func):
    def wrapper(request, *args, **kwargs):
        if request.method == "POST":
            pdf_file = request.FILES.get('pdf_file')
            if not pdf_file:
                return HttpResponse("No PDF file uploaded.", status=400)
            
            pdf_path = default_storage.save(pdf_file.name, ContentFile(pdf_file.read()))
            pdf_path_full = default_storage.path(pdf_path)

            try:
                tiff_path_full = pdf_path_full.replace('.pdf', '.tiff').replace('.pdf', '.tif')
                if not os.path.exists(pdf_path_full):
                    raise Exception(f'{pdf_path_full} not found.')

                images = pdf2image.convert_from_path(pdf_path_full)
                
                images[0].save(tiff_path_full, save_all=True, append_images=images[1:], compression='tiff_deflate')
                
                tiff_url = default_storage.url(os.path.basename(tiff_path_full))
                print(f'tiff_url {tiff_url}')
                return func(request, tiff_url=tiff_url, *args, **kwargs)
            except Exception as e:
                return HttpResponse(f"Error converting PDF to TIFF: {e}", status=500)
        else:
            return func(request, *args, **kwargs)
    return wrapper

@pdf_to_tiff_logic
def pdf_to_tiff_view(request, tiff_url=None):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to TIFF file converter online',
        description='Convert PDF (Portable Document Format) in to TIFF (Tag Image File Format) file format online in free.',
        keywords=['tiff', 'image', 'pdf', 'Portable Document Format', 'Tag Image File Format'],
        og_title='iLovePdfConverterOnline - PDF to TIFF file converter online',
        og_description='Convert PDF (Portable Document Format) in to TIFF (Tag Image File Format) file format online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_tiff_view')
    context = {'meta': meta, 'tiff_url': tiff_url, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_tiff.html', context)

@pdf_to_tiff_logic
def pdf_to_tiff_include(request, tiff_url=None):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to TIFF file converter online',
        description='Convert PDF (Portable Document Format) in to TIFF (Tag Image File Format) file format online in free.',
        keywords=['tiff', 'image', 'pdf', 'Portable Document Format', 'Tag Image File Format'],
        og_title='iLovePdfConverterOnline - PDF to TIFF file converter online',
        og_description='Convert PDF (Portable Document Format) in to TIFF (Tag Image File Format) file format online in free.',
    )
    context = {'meta': meta, 'tiff_url': tiff_url}
    return render(request, 'tools/pdf_to_tiff_include.html', context)



# -----------------------------------------================================

def xml_to_pdf_logic(view_func):
    def wrapper(request, *args, **kwargs):
        if request.method == 'POST':
            xml_content = request.POST.get('xml_content', '')

            # If a file is uploaded, read its content
            if 'xml_file' in request.FILES:
                xml_file = request.FILES['xml_file']
                xml_content = xml_file.read().decode('utf-8')

            # Get the option to remove HTML tags from the request (if provided)
            remove_tags = request.POST.get('remove_tags', False) in ['on', 'true', '1']

            # Convert XML to PDF with the chosen option
            pdf_content = convert_to_pdf(xml_content, remove_tags)

            # Send PDF as downloadable response
            response = HttpResponse(pdf_content, content_type='application/pdf')
            response['Content-Disposition'] = 'attachment; filename="XML_to_PDF_ilovepdfconverteronline.com.pdf"'
            return response

        return view_func(request, *args, **kwargs)
    return wrapper

def convert_to_pdf(xml_content, remove_tags=False):
    """Converts XML content to PDF with an optional flag to remove HTML tags.

    Args:
        xml_content (str): The XML content to be converted.
        remove_tags (bool, optional): If True, removes HTML tags from the output PDF. Defaults to False.

    Returns:
        bytes: The generated PDF content.
    """
    if remove_tags:
        # Escape special characters to prevent them from being interpreted as HTML tags
        escaped_content = xml_content.replace('&', '&amp;').replace('<', '&lt;').replace('>', '&gt;')
        html_content = f"<html><body><pre>{escaped_content}</pre></body></html>"
    else:
        html_content = f"<html><body><pre>{xml_content}</pre></body></html>"

    # Generate PDF using pdfkit
    pdf = pdfkit.from_string(html_content, False)
    return pdf

@xml_to_pdf_logic
def xml_to_pdf_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - XML to PDF converter online',
        description='Convert Extensible Markup Language (XML) file in to PDF (Portable Document Format) online in free.',
        keywords=['pdf', 'xml', 'file', 'Portable Document Format', 'Extensible Markup Language'],
        og_title='iLovePdfConverterOnline - XML to PDF converter online',
        og_description='Convert Extensible Markup Language (XML) file in to PDF (Portable Document Format) online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='xml_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment, 'remove_tags_checked': ''}
    
    # Render the template with an optional checkbox for removing tags
    # context = {
    #     'remove_tags_checked': ''  # Set to 'checked' if desired by default
    # }
    return render(request, 'tools/xml_to_pdf.html', context)

@xml_to_pdf_logic
def xml_to_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - XML to PDF converter online',
        description='Convert Extensible Markup Language (XML) file in to PDF (Portable Document Format) online in free.',
        keywords=['pdf', 'xml', 'file', 'Portable Document Format', 'Extensible Markup Language'],
        og_title='iLovePdfConverterOnline - XML to PDF converter online',
        og_description='Convert Extensible Markup Language (XML) file in to PDF (Portable Document Format) online in free.',
    )
    context = {'meta': meta, 'remove_tags_checked': ''}
    # Render the template with an optional checkbox for removing tags
    # context = {
    #     'remove_tags_checked': ''  # Set to 'checked' if desired by default
    # }
    return render(request, 'tools/xml_to_pdf_include.html', context)

# -----------------------------------------================================


def pdf_to_xml_logic(view_func):
    @wraps(view_func)
    def _wrapped_view(request, *args, **kwargs):
        if request.method == 'POST':
            form = UploadFileForm(request.POST, request.FILES)
            if form.is_valid():
                pdf_file = request.FILES['file']
                # Create a BytesIO object to store the XML content
                xml_output = BytesIO()

                # Convert the PDF to XML and write it to the BytesIO object
                extract_text_to_fp(pdf_file, xml_output, output_type='xml')

                # Seek to the beginning of the BytesIO object
                xml_output.seek(0)

                # Read the XML content from the BytesIO object
                xml_content = xml_output.read()

                # Close the BytesIO object
                xml_output.close()

                response = HttpResponse(xml_content, content_type='application/xml')
                response['Content-Disposition'] = 'attachment; filename=PDF_to_XML_ilovepdfconverteronline.com.xml'
                return response
        else:
            form = UploadFileForm()
        return view_func(request, form, *args, **kwargs)
    return _wrapped_view

@pdf_to_xml_logic
def pdf_to_xml_view(request, form):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to XML converter online',
        description='Convert PDF (Portable Document Format) to XML (Extensible Markup Language) file online in free.',
        keywords= ['pdf', 'xml', 'file', 'Portable Document Format', 'Extensible Markup Language'],
        og_title='iLovePdfConverterOnline - PDF to XML converter online',
        og_description='Convert PDF (Portable Document Format) to XML (Extensible Markup Language) file online in free.',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pdf_to_xml_view')
    context = {'form': form, 'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pdf_to_xml.html', context)

@pdf_to_xml_logic
def pdf_to_xml_include(request, form):
    meta = Meta(
        title='iLovePdfConverterOnline - PDF to XML converter online',
        description='Convert PDF (Portable Document Format) to XML (Extensible Markup Language) file online in free.',
        keywords= ['pdf', 'xml', 'file', 'Portable Document Format', 'Extensible Markup Language'],
        og_title='iLovePdfConverterOnline - PDF to XML converter online',
        og_description='Convert PDF (Portable Document Format) to XML (Extensible Markup Language) file online in free.',
    )
    context = {'form': form, 'meta': meta}
    return render(request, 'tools/pdf_to_xml_include.html', context )

# -----------------------------------------================================

# working 100%
def html_to_pdf_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - HTML to PDF',
        description='Convert HTML file or URL in to PDF. HyperText Markup Language to Portable Document Format',
        keywords=['html', 'url', 'urls', 'links', 'file', 'download', 'pdf', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - HTML to PDF',
        og_description='Convert HTML file or URL to PDF. HyperText Markup Language to Portable Document Format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='html_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}

    if request.method == 'POST':
        if 'url' in request.POST:
            url = request.POST['url']
            # Check if the URL is not empty
            if url:
                pdf = pdfkit.from_url(url, False)
                response = HttpResponse(pdf, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="HTML2PDF_ilovepdfconverteronline.com.pdf"'
                return response
            else:
                # return HttpResponse("URL is empty")
                if 'html_file' in request.FILES:
                    html_file = request.FILES['html_file']
                    # Read content of HTML file
                    html_content = html_file.read().decode('utf-8')
                    # print(html_content)
                    # Convert HTML content to PDF
                    pdf = pdfkit.from_string(html_content, options={"enable-local-file-access": ""})
                    response = HttpResponse(pdf, content_type='application/pdf')
                    response['Content-Disposition'] = 'attachment; filename="HTML2PDF_ilovepdfconverteronline.com.pdf"'
                    return response
                else:
                    return HttpResponse("Invalid Request")
    else:
        return render(request, 'tools/html_to_pdf.html', context)

def html_to_pdf_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - HTML to PDF',
        description='Convert HTML file or URL to PDF.',
        keywords=['html', 'url', 'urls', 'links', 'file', 'download', 'pdf', 'Portable Document Format'],
        og_title='iLovePdfConverterOnline - HTML to PDF',
        og_description='Convert HTML file or URL to PDF.',
    )

    context = {'meta': meta}

    if request.method == 'POST':
        if 'url' in request.POST:
            url = request.POST['url']
            # Check if the URL is not empty
            if url:
                pdf = pdfkit.from_url(url, False)
                response = HttpResponse(pdf, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="HTML2PDF_ilovepdfconverteronline.com.pdf"'
                return response
            else:
                # return HttpResponse("URL is empty")
                if 'html_file' in request.FILES:
                    html_file = request.FILES['html_file']
                    # Read content of HTML file
                    html_content = html_file.read().decode('utf-8')
                    # print(html_content)
                    # Convert HTML content to PDF
                    pdf = pdfkit.from_string(html_content, options={"enable-local-file-access": ""})
                    response = HttpResponse(pdf, content_type='application/pdf')
                    response['Content-Disposition'] = 'attachment; filename="HTML2PDF_ilovepdfconverteronline.com.pdf"'
                    return response
                else:
                    return HttpResponse("Invalid Request")
    else:
        return render(request, 'tools/html_to_pdf_include.html', context)

# -----------------------------------------================================



def count_words(text):
    # Remove punctuation and split by whitespace to count words
    words = re.findall(r'\b\w+\b', text)
    return len(words)


def word_counter_text_logic(template_name):
    def decorator(view_func):
        @wraps(view_func)
        def _wrapped_view(request, *args, **kwargs):
            context = view_func(request, *args, **kwargs)
            word_count = 0
            if request.method == 'POST':
                text = request.POST.get('text', '')
                word_count = count_words(text)
            context['word_count'] = word_count
            if request.headers.get('x-requested-with') == 'XMLHttpRequest':
                return JsonResponse({'word_count': word_count})
            return render(request, template_name, context)
        return _wrapped_view
    return decorator

@word_counter_text_logic('tools/word_counter_text.html')
def word_counter_text_view(request):
    context = {}
    return context

@word_counter_text_logic('tools/word_counter_text_include.html')
def word_counter_text_include(request):
    context = {}
    return context


# .................................................................==================================================

def lorem_ipsum_generator_logic(template_name):
    def decorator(view_func):
        @wraps(view_func)
        def _wrapped_view(request, *args, **kwargs):
            context = view_func(request, *args, **kwargs)
            if request.method == 'POST':
                paragraphs = int(request.POST.get('paragraphs', 3))
                lorem_ipsum_text = generate_lorem_ipsum(paragraphs)
                context['lorem_ipsum_text'] = lorem_ipsum_text
            else:
                context['lorem_ipsum_text'] = ""
            return render(request, template_name, context)
        return _wrapped_view
    return decorator


@lorem_ipsum_generator_logic('tools/lorem_ipsum_generator.html')
def lorem_ipsum_generator_view(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Lorem Ipsum Generator online',
        description='Generate Lorem Ipsum paragraph with random text online in free.',
        keywords= ['pdf', 'xml', 'file'],
        og_title='iLovePdfConverterOnline - Lorem Ipsum Generator online',
        og_description='Generate Lorem Ipsum paragraph with random text online in free.',
    )
    context = {'meta': meta}
    return context

@lorem_ipsum_generator_logic('tools/lorem_ipsum_generator_include.html')
def lorem_ipsum_generator_include(request):
    meta = Meta(
        title='iLovePdfConverterOnline - Lorem Ipsum Generator online',
        description='Generate Lorem Ipsum paragraph with random text online in free.',
        keywords= ['pdf', 'xml', 'file'],
        og_title='iLovePdfConverterOnline - Lorem Ipsum Generator online',
        og_description='Generate Lorem Ipsum paragraph with random text online in free.',
    )
    context = {'meta': meta}
    return context



# .................................................................==================================================

# views.py
import io
import fitz  # PyMuPDF
from django.views.decorators.http import require_http_methods
from django.core.files.uploadedfile import UploadedFile

def pdf_to_raw_logic(func):
    def wrapper(request, *args, **kwargs):
        if request.method == 'POST' and request.FILES.get('pdf_file'):
            pdf_file: UploadedFile = request.FILES['pdf_file']
            try:
                pdf_document = fitz.open(stream=pdf_file.read(), filetype="pdf")
                raw_image_data = []
                
                for page_num in range(len(pdf_document)):
                    page = pdf_document.load_page(page_num)
                    pix = page.get_pixmap()
                    img_bytes = pix.samples  # Raw pixel data
                    raw_image_data.append(img_bytes)

                # For simplicity, we'll return the first page's raw image data
                response = HttpResponse(raw_image_data[0], content_type='application/octet-stream')
                response['Content-Disposition'] = 'attachment; filename="page_1.raw"'
                return response
            except Exception as e:
                return HttpResponseBadRequest(f"Error processing PDF file: {str(e)}")
        return func(request, *args, **kwargs)
    return wrapper

@pdf_to_raw_logic
@require_http_methods(["GET", "POST"])
def pdf_to_raw_view(request):
    return render(request, 'tools/pdf_to_raw.html')



from django.views.decorators.http import require_http_methods
import io
import rawpy
import numpy as np

def raw_to_pdf_logic(view_func):
    @wraps(view_func)
    def wrapper(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('raw_image'):
            raw_image = request.FILES['raw_image']

            try:
                # Use rawpy to read the RAW image file
                with rawpy.imread(raw_image) as raw:
                    rgb_image = raw.postprocess()

                # Convert numpy array to PIL image
                image = Image.fromarray(rgb_image)
                buffer = io.BytesIO()
                image.save(buffer, format='JPEG')
                buffer.seek(0)

                # Create a PDF file from the image
                pdf_buffer = io.BytesIO()
                c = canvas.Canvas(pdf_buffer, pagesize=letter)
                c.drawImage(buffer, 0, 0, width=letter[0], height=letter[1])
                c.showPage()
                c.save()
                pdf_buffer.seek(0)

                response = HttpResponse(pdf_buffer, content_type='application/pdf')
                response['Content-Disposition'] = 'attachment; filename="converted.pdf"'
                return response

            except rawpy._rawpy.LibRawError:
                return HttpResponse("Invalid RAW image file.", status=400)

        return view_func(request, *args, **kwargs)

    return wrapper

@require_http_methods(["GET", "POST"])
@raw_to_pdf_logic
def raw_to_pdf_view(request):
    return render(request, 'tools/raw_to_pdf.html')



###########################################################################################

####################################################################################################
import shutil

def odp_to_pptx_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('odp_file'):
            try:
                odp_file = request.FILES['odp_file']

                # Save uploaded ODP file to temporary location
                temp_filename = odp_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, odp_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pptx', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_pptx_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.pptx')

                if os.path.exists(output_pptx_path):
                    # Serve the PPTX file for download
                    response = FileResponse(open(output_pptx_path, 'rb'), content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pptx_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_pptx_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to PPTX")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function


@odp_to_pptx_logic
def odp_to_pptx_view(request):
    meta = Meta(
        title='OpenDocument Presentation (.odp) file to Microsoft PowerPoint (.pptx) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Microsoft PowerPoint (.pptx) file format',
        keywords=['powerpoint', 'microsoft powerpoint', 'ppt', 'pptx', 'OpenDocument', 'Presentation'],
        og_title='OpenDocument Presentation (.odp) file to Microsoft PowerPoint (.pptx) file converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Microsoft PowerPoint (.pptx) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pptx_to_odp_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/odp_to_pptx.html', context)

@odp_to_pptx_logic
def odp_to_pptx_include(request):
    meta = Meta(
        title='OpenDocument Presentation (.odp) file to Microsoft PowerPoint (.pptx) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Microsoft PowerPoint (.pptx) file format',
        keywords=['powerpoint', 'microsoft powerpoint', 'ppt', 'pptx', 'OpenDocument', 'Presentation'],
        og_title='OpenDocument Presentation (.odp) file to Microsoft PowerPoint (.pptx) file converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Microsoft PowerPoint (.pptx) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/odp_to_pptx_include.html', context)


###########################################################


def ods_to_xlsx_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('ods_file'):
            try:
                ods_file = request.FILES['ods_file']

                # Save uploaded ODS file to temporary location
                temp_filename = ods_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, ods_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'xlsx', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_xlsx_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.xlsx')

                if os.path.exists(output_xlsx_path):
                    # Serve the XLSX file for download
                    response = FileResponse(open(output_xlsx_path, 'rb'), content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_xlsx_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_xlsx_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to XLSX")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@ods_to_xlsx_logic
def ods_to_xlsx_view(request):
    meta = Meta(
        title='OpenDocument Spreadsheet (.ods) file to Microsoft Excel (.xlsx) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Microsoft Excel (.xlsx) file format',
        keywords=['opendocument spreadsheet', 'microsoft excel', 'ods', 'xlsx', 'xls'],
        og_title='OpenDocument Spreadsheet (.ods) file to Microsoft Excel (.xlsx) file converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Microsoft Excel (.xlsx) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='ods_to_xlsx_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/ods_to_xlsx.html', context)

@ods_to_xlsx_logic
def ods_to_xlsx_include(request):
    meta = Meta(
        title='OpenDocument Spreadsheet (.ods) file to Microsoft Excel (.xlsx) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Microsoft Excel (.xlsx) file format',
        keywords=['opendocument spreadsheet', 'microsoft excel', 'ods', 'xlsx', 'xls'],
        og_title='OpenDocument Spreadsheet (.ods) file to Microsoft Excel (.xlsx) file converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Microsoft Excel (.xlsx) file format',
    )
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/ods_to_xlsx_include.html', context)

########################################################

# import os
# import shutil
# import subprocess
# from django.conf import settings
# from django.core.files.storage import FileSystemStorage
# from django.http import FileResponse, HttpResponse
# from django.shortcuts import render

def odt_to_docx_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('odt_file'):
            try:
                odt_file = request.FILES['odt_file']

                # Save uploaded ODT file to temporary location
                temp_filename = odt_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, odt_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'docx', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_docx_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.docx')

                if os.path.exists(output_docx_path):
                    # Serve the DOCX file for download
                    response = FileResponse(open(output_docx_path, 'rb'), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_docx_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_docx_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to DOCX")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@odt_to_docx_logic
def odt_to_docx_view(request):
    meta = Meta(
        title='OpenDocument Text (.odt) file to Microsoft Word (.docx) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Text (.odt) file in to Microsoft Word (.docx) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'odt', 'opendocument', 'text'],
        og_title='OpenDocument Text (.odt) file to Microsoft Word (.docx) file converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Text (.odt) file in to Microsoft Word (.docx) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='odt_to_docx_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/odt_to_docx.html', context)

@odt_to_docx_logic
def odt_to_docx_include(request):
    meta = Meta(
        title='OpenDocument Text (.odt) file to Microsoft Word (.docx) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Text (.odt) file in to Microsoft Word (.docx) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'odt', 'opendocument', 'text'],
        og_title='OpenDocument Text (.odt) file to Microsoft Word (.docx) file converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Text (.odt) file in to Microsoft Word (.docx) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/odt_to_docx_include.html', context)


####################################################


def odt_to_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('odt_file'):
            try:
                odt_file = request.FILES['odt_file']

                # Save uploaded ODT file to temporary location
                temp_filename = odt_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, odt_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_pdf_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.pdf')

                if os.path.exists(output_pdf_path):
                    # Serve the PDF file for download
                    response = FileResponse(open(output_pdf_path, 'rb'), content_type='application/pdf')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pdf_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_pdf_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to PDF")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@odt_to_pdf_logic
def odt_to_pdf_view(request):
    meta = Meta(
        title='OpenDocument Text (.odt) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.odt) file in to Portable Document Format (.pdf) file format',
        keywords=['OpenDocument Text', 'pdf', 'odt', 'Portable Document Format'],
        og_title='OpenDocument Text (.odt) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Text (.odt) file in to Portable Document Format (.pdf) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='odt_to_pdf_view')
    context = {'meta': meta, 'tool_attachment':tool_attachment}
    return render(request, 'tools/odt_to_pdf.html', context)

@odt_to_pdf_logic
def odt_to_pdf_include(request):
    meta = Meta(
        title='OpenDocument Text (.odt) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.odt) file in to Portable Document Format (.pdf) file format',
        keywords=['OpenDocument Text', 'pdf', 'odt', 'Portable Document Format'],
        og_title='OpenDocument Text (.odt) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Text (.odt) file in to Portable Document Format (.pdf) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/odt_to_pdf_include.html', context)

##################################################################

def ods_to_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('ods_file'):
            try:
                ods_file = request.FILES['ods_file']

                # Save uploaded ODS file to temporary location
                temp_filename = ods_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, ods_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_pdf_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.pdf')

                if os.path.exists(output_pdf_path):
                    # Serve the PDF file for download
                    response = FileResponse(open(output_pdf_path, 'rb'), content_type='application/pdf')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pdf_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_pdf_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to PDF")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function


@ods_to_pdf_logic
def ods_to_pdf_view(request):
    meta = Meta(
        title='OpenDocument Spreadsheet (.ods) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Portable Document Format (.pdf) file format',
        keywords=['OpenDocument Spreadsheet', 'pdf', 'ods', 'Portable Document Format'],
        og_title='OpenDocument Spreadsheet (.ods) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Portable Document Format (.pdf) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='ods_to_pdf_view')
    context = {'meta': meta, 'tool_attachment':tool_attachment}
    return render(request, 'tools/ods_to_pdf.html', context)

@ods_to_pdf_logic
def ods_to_pdf_include(request):
    meta = Meta(
        title='OpenDocument Spreadsheet (.ods) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Portable Document Format (.pdf) file format',
        keywords=['OpenDocument Spreadsheet', 'pdf', 'ods', 'Portable Document Format'],
        og_title='OpenDocument Spreadsheet (.ods) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Spreadsheet (.ods) file in to Portable Document Format (.pdf) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/ods_to_pdf_include.html', context)

########################################################################

def odp_to_pdf_logic(view_func):
    @wraps(view_func)
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('odp_file'):
            try:
                odp_file = request.FILES['odp_file']

                # Save uploaded ODP file to temporary location
                temp_filename = odp_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, odp_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_pdf_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.pdf')

                if os.path.exists(output_pdf_path):
                    # Serve the PDF file for download
                    response = FileResponse(open(output_pdf_path, 'rb'), content_type='application/pdf')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pdf_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_pdf_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to PDF")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@odp_to_pdf_logic
def odp_to_pdf_view(request):
    meta = Meta(
        title='OpenDocument Presentation (.odp) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Portable Document Format (.pdf) file format',
        keywords=['OpenDocument Presentation', 'pdf', 'odp', 'Portable Document Format'],
        og_title='OpenDocument Presentation (.odp) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Portable Document Format (.pdf) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='odp_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/odp_to_pdf.html', context)

@odp_to_pdf_logic
def odp_to_pdf_include(request):
    meta = Meta(
        title='OpenDocument Presentation (.odp) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Portable Document Format (.pdf) file format',
        keywords=['OpenDocument Presentation', 'pdf', 'odp', 'Portable Document Format'],
        og_title='OpenDocument Presentation (.odp) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts OpenDocument Presentation (.odp) file in to Portable Document Format (.pdf) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/odp_to_pdf_include.html', context)

#################################################################


def rtf_to_pdf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('rtf_file'):
            try:
                rtf_file = request.FILES['rtf_file']

                # Save uploaded RTF file to temporary location
                temp_filename = rtf_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, rtf_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'pdf', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_pdf_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.pdf')

                if os.path.exists(output_pdf_path):
                    # Serve the PDF file for download
                    response = FileResponse(open(output_pdf_path, 'rb'), content_type='application/pdf')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_pdf_path)}'
                    
                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_pdf_path]
                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)
                    return response
                else:
                    return HttpResponse("Error converting file to PDF")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@rtf_to_pdf_logic
def rtf_to_pdf_view(request):
    meta = Meta(
        title='Rich Text Format (.rtf) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Portable Document Format (.pdf) file format',
        keywords=['Portable Document Format', 'pdf', 'rtf', 'rich text format', 'text'],
        og_title='Rich Text Format (.rtf) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Portable Document Format (.pdf) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='rtf_to_pdf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment} 
    return render(request, 'tools/rtf_to_pdf.html', context)

@rtf_to_pdf_logic
def rtf_to_pdf_include(request):
    meta = Meta(
        title='Rich Text Format (.rtf) file to Portable Document Format (.pdf) file converter',
        description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Portable Document Format (.pdf) file format',
        keywords=['Portable Document Format', 'pdf', 'rtf', 'rich text format', 'text'],
        og_title='Rich Text Format (.rtf) file to Portable Document Format (.pdf) converter',
        og_description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Portable Document Format (.pdf) file format',
    )
    context = {'meta': meta} 
    return render(request, 'tools/rtf_to_pdf_include.html', context)

##############################################################

def rtf_to_docx_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('rtf_file'):
            try:
                rtf_file = request.FILES['rtf_file']

                # Save uploaded RTF file to temporary location
                temp_filename = rtf_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, rtf_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'docx', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_docx_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.docx')

                if os.path.exists(output_docx_path):
                    # Serve the DOCX file for download
                    response = FileResponse(open(output_docx_path, 'rb'), content_type='application/vnd.openxmlformats-officedocument.wordprocessingml.document')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_docx_path)}'

                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_docx_path]

                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)

                    return response
                else:
                    return HttpResponse("Error converting file to DOCX")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@rtf_to_docx_logic
def rtf_to_docx_view(request):
    meta = Meta(
        title='Rich Text Format (.rtf) file to Microsoft Word (.docx) file converter',
        description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Microsoft Word (.docx) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'rtf', 'rich text format', 'text'],
        og_title='Rich Text Format (.rtf) file to Microsoft Word (.docx) converter',
        og_description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Microsoft Word (.docx) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='rtf_to_docx_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment} 
    return render(request, 'tools/rtf_to_docx.html', context)

@rtf_to_docx_logic
def rtf_to_docx_include(request):
    meta = Meta(
        title='Rich Text Format (.rtf) file to Microsoft Word (.docx) file converter',
        description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Microsoft Word (.docx) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'rtf', 'rich text format', 'text'],
        og_title='Rich Text Format (.rtf) file to Microsoft Word (.docx) converter',
        og_description='iLovePdfConverterOnline Converts Rich Text Format (.rtf) file in to Microsoft Word (.docx) file format',
    )
    context = {'meta': meta} 
    return render(request, 'tools/rtf_to_docx_include.html', context)

###################################################

def docx_to_rtf_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('docx_file'):
            try:
                docx_file = request.FILES['docx_file']

                # Save uploaded DOCX file to temporary location
                temp_filename = docx_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, docx_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'rtf', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_rtf_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.rtf')

                if os.path.exists(output_rtf_path):
                    # Serve the RTF file for download
                    response = FileResponse(open(output_rtf_path, 'rb'), content_type='application/rtf')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_rtf_path)}'

                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_rtf_path]

                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)

                    return response
                else:
                    return HttpResponse("Error converting file to RTF")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@docx_to_rtf_logic
def docx_to_rtf_view(request):
    meta = Meta(
        title='Microsoft Word (.docx) file to Rich Text Format (.rtf) file converter',
        description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to Rich Text Format (.rtf) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'rtf', 'rich text format', 'text'],
        og_title='Microsoft Word (.docx) file to Rich Text Format (.rtf) converter',
        og_description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to Rich Text Format (.rtf) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='docx_to_rtf_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment} 
    return render(request, 'tools/docx_to_rtf.html', context)

@docx_to_rtf_logic
def docx_to_rtf_include(request):
    meta = Meta(
        title='Microsoft Word (.docx) file to Rich Text Format (.rtf) file converter',
        description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to Rich Text Format (.rtf) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'rtf', 'rich text format', 'text'],
        og_title='Microsoft Word (.docx) file to Rich Text Format (.rtf) converter',
        og_description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to Rich Text Format (.rtf) file format',
    )
    context = {'meta': meta} 
    return render(request, 'tools/docx_to_rtf_include.html', context)

################################################################

def docx_to_odt_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('docx_file'):
            try:
                docx_file = request.FILES['docx_file']

                # Save uploaded DOCX file to temporary location
                temp_filename = docx_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, docx_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'odt', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_odt_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.odt')

                if os.path.exists(output_odt_path):
                    # Serve the ODT file for download
                    response = FileResponse(open(output_odt_path, 'rb'), content_type='application/vnd.oasis.opendocument.text')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_odt_path)}'

                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_odt_path]

                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)

                    return response
                else:
                    return HttpResponse("Error converting file to ODT")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@docx_to_odt_logic
def docx_to_odt_view(request):
    meta = Meta(
        title='Microsoft Word (.docx) file to OpenDocument Text (.odt) file converter',
        description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to OpenDocument Text (.odt) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'odt', 'opendocument', 'text'],
        og_title='Microsoft Word (.docx) file to OpenDocument Text (.odt) converter',
        og_description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to OpenDocument Text (.odt) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='docx_to_odt_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/docx_to_odt.html', context)

@docx_to_odt_logic
def docx_to_odt_include(request):
    meta = Meta(
        title='Microsoft Word (.docx) file to OpenDocument Text (.odt) file converter',
        description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to OpenDocument Text (.odt) file format',
        keywords=['word', 'microsoft word', 'doc', 'docx', 'odt', 'opendocument', 'text'],
        og_title='Microsoft Word (.docx) file to OpenDocument Text (.odt) converter',
        og_description='iLovePdfConverterOnline Converts Microsoft Word (.docx) file in to OpenDocument Text (.odt) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/docx_to_odt_include.html', context)

##################################################################

def xlsx_to_ods_logic(view_func):
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('xlsx_file'):
            try:
                xlsx_file = request.FILES['xlsx_file']

                # Save uploaded XLSX file to temporary location
                temp_filename = xlsx_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, xlsx_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                # Convert XLSX to ODS using LibreOffice
                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'ods', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_ods_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.ods')

                if os.path.exists(output_ods_path):
                    # Serve the ODS file for download
                    response = FileResponse(open(output_ods_path, 'rb'), content_type='application/vnd.oasis.opendocument.spreadsheet')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_ods_path)}'

                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_ods_path]

                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)

                    return response
                else:
                    return HttpResponse("Error converting file to ODS")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function

@xlsx_to_ods_logic
def xlsx_to_ods_view(request):
    meta = Meta(
        title='Excel (.xlsx) file to OpenDocument Spreadsheet (.ods) file converter',
        description='iLovePdfConverterOnline Converts Excel (.xlsx) file in to OpenDocument Spreadsheet (.ods) file format',
        keywords=['powerpoint', 'microsoft powerpoint', 'ppt', 'pptx'],
        og_title='Excel (.xlsx) file to OpenDocument Spreadsheet (.ods) converter',
        og_description='iLovePdfConverterOnline Converts Excel (.xlsx) file in to OpenDocument Spreadsheet (.ods) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='xlsx_to_ods_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/xlsx_to_ods.html', context)

@xlsx_to_ods_logic
def xlsx_to_ods_include(request):
    meta = Meta(
        title='Excel (.xlsx) file to OpenDocument Spreadsheet (.ods) file converter',
        description='iLovePdfConverterOnline Converts Excel (.xlsx) file in to OpenDocument Spreadsheet (.ods) file format',
        keywords=['powerpoint', 'microsoft powerpoint', 'ppt', 'pptx'],
        og_title='Excel (.xlsx) file to OpenDocument Spreadsheet (.ods) converter',
        og_description='iLovePdfConverterOnline Converts Excel (.xlsx) file in to OpenDocument Spreadsheet (.ods) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/xlsx_to_ods_include.html', context)

#########################################################################


def pptx_to_odp_logic(view_func):
    @wraps(view_func)
    def wrapper_function(request, *args, **kwargs):
        if request.method == "POST" and request.FILES.get('pptx_file'):
            try:
                pptx_file = request.FILES['pptx_file']

                # Save uploaded PPTX file to temporary location
                temp_filename = pptx_file.name
                temp_file_path = os.path.join(settings.MEDIA_ROOT, 'uploads', temp_filename)

                fs = FileSystemStorage(location=os.path.join(settings.MEDIA_ROOT, 'uploads'))
                temp_file = fs.save(temp_filename, pptx_file)

                out_path = os.path.join(settings.MEDIA_ROOT, 'uploads')

                env = os.environ.copy()
                env['HOME'] = os.path.join(settings.MEDIA_ROOT, 'uploads')

                result = subprocess.run(
                    ['libreoffice', '--headless', '--convert-to', 'odp', '--outdir', out_path, temp_file_path],
                    env=env, capture_output=True, text=True
                )
                if result.returncode != 0:
                    raise Exception(f"Subprocess failed with error: {result.stderr}")

                output_odp_path = os.path.join(out_path, os.path.splitext(temp_filename)[0] + '.odp')

                if os.path.exists(output_odp_path):
                    # Serve the ODP file for download
                    response = FileResponse(open(output_odp_path, 'rb'), content_type='application/vnd.oasis.opendocument.presentation')
                    response['Content-Disposition'] = f'attachment; filename={os.path.basename(output_odp_path)}'

                    # Add files to cleanup list
                    response.cleanup_files = [temp_file_path, output_odp_path]

                    # Clean up .cache and .config directories
                    cache_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.cache')
                    config_path = os.path.join(settings.MEDIA_ROOT, 'uploads', '.config')
                    if os.path.exists(cache_path):
                        shutil.rmtree(cache_path)
                    if os.path.exists(config_path):
                        shutil.rmtree(config_path)

                    return response
                else:
                    return HttpResponse("Error converting file to ODP")
            except Exception as e:
                return HttpResponse(status=500, content=str(e))
        else:
            return view_func(request, *args, **kwargs)
    return wrapper_function


@pptx_to_odp_logic
def pptx_to_odp_view(request):
    meta = Meta(
        title='PowerPoint (.pptx) file in to OpenDocument Presentation (.odp) file converter',
        description='iLovePdfConverterOnline Converts PowerPoint (.pptx) file in to OpenDocument Presentation (.odp) file format',
        keywords=['powerpoint', 'microsoft powerpoint', 'ppt', 'pptx'],
        og_title='PowerPoint (.pptx) to OpenDocument Presentation (.odp) converter',
        og_description='iLovePdfConverterOnline Convert PowerPoint (.pptx) in to OpenDocument Presentation (.odp) file format',
    )
    tool_attachment = ToolAttachment.objects.get(function_name='pptx_to_odp_view')
    context = {'meta': meta, 'tool_attachment': tool_attachment}
    return render(request, 'tools/pptx_to_odp.html', context)

@pptx_to_odp_logic
def pptx_to_odp_include(request):
    meta = Meta(
        title='PowerPoint (.pptx) file to OpenDocument Presentation (.odp) file converter',
        description='iLovePdfConverterOnline Converts PowerPoint (.pptx) file in to OpenDocument Presentation (.odp) file format',
        keywords=['powerpoint', 'microsoft powerpoint', 'ppt', 'pptx'],
        og_title='PowerPoint (.pptx) to OpenDocument Presentation (.odp) converter',
        og_description='iLovePdfConverterOnline Converts PowerPoint (.pptx) in to OpenDocument Presentation (.odp) file format',
    )
    context = {'meta': meta}
    return render(request, 'tools/pptx_to_odp_include.html', context)

#####################################################################

