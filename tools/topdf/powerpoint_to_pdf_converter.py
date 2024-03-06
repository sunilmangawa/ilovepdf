from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
import os

def convert_ppt_to_pdf(ppt_file):
    presentation = Presentation(ppt_file)
    pdf_path = os.path.splitext(ppt_file)[0] + '.pdf'

    pdf = canvas.Canvas(pdf_path, pagesize=letter)
    width, height = letter

    for slide in presentation.slides:
        background = slide.background
        if background.fill.background():
            bg_image = background.fill.background().blip
            if bg_image:
                image_path = os.path.join('media', 'temp_bg.png')  # Temporary file for background image
                bg_image.save(image_path)
                pdf.drawImage(image_path, 0, 0, width=width, height=height, preserveAspectRatio=True)
                os.remove(image_path)  # Remove the temporary background image

        # Draw the content of the slide
        for shape in slide.shapes:
            if shape.has_text_frame:
                text_frame = shape.text_frame
                left = shape.left if shape.left is not None else 0
                top = shape.top if shape.top is not None else 0
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        font_size = run.font.size if run.font.size is not None else 12  # Default font size
                        pdf.drawString(left, top - font_size, run.text)

        pdf.showPage()

    pdf.save()
    return pdf_path


# from pptx import Presentation
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# import os

# def convert_ppt_to_pdf(ppt_file):
#     presentation = Presentation(ppt_file)
#     pdf_path = os.path.splitext(ppt_file)[0] + '.pdf'

#     pdf = canvas.Canvas(pdf_path, pagesize=letter)
#     width, height = letter

#     for slide in presentation.slides:
#         background = slide.background
#         if background.fill.background():
#             bg_image = background.fill.background().blip
#             if bg_image:
#                 image_path = os.path.join('media', 'temp_bg.png')  # Temporary file for background image
#                 bg_image.save(image_path)
#                 pdf.drawImage(image_path, 0, 0, width=width, height=height, preserveAspectRatio=True)
#                 os.remove(image_path)  # Remove the temporary background image

#         # Draw the content of the slide
#         for shape in slide.shapes:
#             if shape.has_text_frame:
#                 text_frame = shape.text_frame
#                 for paragraph in text_frame.paragraphs:
#                     for run in paragraph.runs:
#                         pdf.drawString(shape.left, shape.top - run.font.size, run.text)

#         pdf.showPage()

#     pdf.save()
#     return pdf_path


# from pptx import Presentation
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# import os

# def convert_ppt_to_pdf(ppt_file):
#     presentation = Presentation(ppt_file)
#     pdf_path = os.path.splitext(ppt_file)[0] + '.pdf'

#     pdf = canvas.Canvas(pdf_path, pagesize=letter)
#     width, height = letter

#     for slide in presentation.slides:
#         background = slide.background
#         if background.fill.background():
#             bg_image = background.fill.background().blip
#             if bg_image:
#                 image_path = os.path.join('media', 'temp_bg.png')  # Temporary file for background image
#                 bg_image.save(image_path)
#                 pdf.drawImage(image_path, 0, 0, width=width, height=height, preserveAspectRatio=True)
#                 os.remove(image_path)  # Remove the temporary background image

#         pdf.showPage()

#     pdf.save()
#     return pdf_path


# from pptx import Presentation
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# import os

# def convert_ppt_to_pdf(ppt_file):
#     presentation = Presentation(ppt_file)
#     pdf_path = os.path.splitext(ppt_file)[0] + '.pdf'

#     pdf = canvas.Canvas(pdf_path, pagesize=letter)
#     for slide in presentation.slides:
#         pdf.drawImage(slide.background.blob, 0, 0, width=letter[0], height=letter[1], preserveAspectRatio=True)
#         pdf.showPage()

#     pdf.save()  # Close the canvas manually

#     return pdf_path


# from pptx import Presentation
# from reportlab.lib.pagesizes import letter
# from reportlab.pdfgen import canvas
# import os

# def convert_ppt_to_pdf(ppt_file):
#     presentation = Presentation(ppt_file)
#     pdf_path = os.path.splitext(ppt_file)[0] + '.pdf'

#     with canvas.Canvas(pdf_path, pagesize=letter) as pdf:
#         for slide in presentation.slides:
#             pdf.drawImage(slide.background.blob, 0, 0, width=letter[0], height=letter[1], preserveAspectRatio=True)
#             pdf.showPage()

#     return pdf_path
